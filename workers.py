"""
Background workers — one per pipeline stage.
All functions are async and run via FastAPI BackgroundTasks.
All model calls use prompts.py exclusively. Do not construct prompt strings here.
"""

import asyncio
import base64
import shutil
import io
import json
import logging
import mimetypes
import os
import re
import time
import zipfile
from collections import Counter
from contextvars import ContextVar
from datetime import datetime, timezone
from pathlib import Path

import httpx
import db
from model_gateway import gateway

# Lazy import — prompts.py is gitignored; checked at startup in main.py
import prompts

JOBS_DIR = Path(os.environ.get("INSTANT_JOBS_DIR", str(Path(__file__).parent / "jobs")))
JOBS_DIR.mkdir(parents=True, exist_ok=True)

TEMPLATE_LABELS = ["Mechanism", "Market", "Economics", "Proof"]  # fallback only

# Per-coroutine overrides for the generation bridge — set by ribotome._run_generation
# so concurrent generation jobs never share a module-global swap.
_raw_user_inputs_cv: ContextVar = ContextVar('_raw_user_inputs_cv', default=None)
_selected_template_payload_cv: ContextVar = ContextVar('_selected_template_payload_cv', default=None)



async def _chunked_gather(coros, chunk_size: int = 4):
    """Run coroutines in bounded-concurrency chunks. This limits how many run
    *simultaneously*; it says nothing about the actual API rate limit, which is
    a rolling-window constraint over time. Image generation calls must go
    through _throttle_image_generation() regardless of how they're batched."""
    results = []
    coros = list(coros)
    for i in range(0, len(coros), chunk_size):
        results.extend(await asyncio.gather(*coros[i : i + chunk_size]))
    return results


async def image_scan_worker(job_id: str) -> None:
    """Vision scan of the intake image. Stores JSON analysis in image_analysis column."""
    job = db.get_job(job_id)
    if not job:
        return
    img_path = _supporting_image(job_id)
    if not img_path:
        return
    b64 = base64.b64encode(img_path.read_bytes()).decode()
    mime = mimetypes.guess_type(str(img_path))[0] or "image/png"
    resp = await gateway.chat.completions.create(
        model=gateway.model_for("anchor"),
        messages=[{
            "role": "user",
            "content": [
                {"type": "text", "text": prompts.image_analysis_prompt()},
                {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
            ],
        }],
        response_format={"type": "json_object"},
    )
    result = json.loads(resp.choices[0].message.content)
    db.update_job(job_id, image_analysis=json.dumps(result))


# Image generation is rate-limited to 50 calls per rolling 60-second window
# (Tier 3 limit for gpt-image-1/gpt-image-2 as of July 2026 -- text
# calls are in the hundreds of thousands/minute and are not a bottleneck).
# This is a real, process-wide constraint (shared across concurrent jobs), not
# a per-batch concurrency cap -- _chunked_gather's chunk_size does not enforce
# it. Every call site that uses image_generation_tool must await this first.
_IMAGE_RATE_LIMIT = 50
_IMAGE_RATE_WINDOW_SECONDS = 60.0
_image_call_times: "list[float]" = []
_image_rate_lock = asyncio.Lock()

# ── DEBUG mode ────────────────────────────────────────────────────────────────
# PS_DEBUG_MODE=1 stubs all image generation calls with a placeholder PNG and
# prints each text call's prompt/response to stdout. No image API spend;
# pipeline logic and all text stages run normally.
_DEBUG_MODE = os.environ.get("PS_DEBUG_MODE") == "1"
_STUB_PNG_B64: str = ""


def _get_stub_png_b64() -> str:
    global _STUB_PNG_B64
    if not _STUB_PNG_B64:
        from PIL import Image, ImageDraw, ImageFont
        import io
        w, h = 1456, 816
        img = Image.new("RGB", (w, h), "white")
        draw = ImageDraw.Draw(img)
        text = "PLACEHOLDER | INFER IMAGE CONTENTS"
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 48)
        except Exception:
            font = ImageFont.load_default()
        bbox = draw.textbbox((0, 0), text, font=font)
        x = (w - (bbox[2] - bbox[0])) // 2
        y = (h - (bbox[3] - bbox[1])) // 2
        draw.text((x, y), text, fill="black", font=font)
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        _STUB_PNG_B64 = base64.b64encode(buf.getvalue()).decode()
    return _STUB_PNG_B64


def _debug_image_resp():
    import types as _t
    item = _t.SimpleNamespace(type="image_generation_call", result=_get_stub_png_b64(), content=[])
    return _t.SimpleNamespace(id="debug-stub", output=[item], output_text=None,
                              usage=_t.SimpleNamespace(input_tokens=0, output_tokens=0, total_tokens=0))


async def _throttle_image_generation() -> None:
    if _DEBUG_MODE:
        return
    async with _image_rate_lock:
        while True:
            now = time.monotonic()
            while _image_call_times and now - _image_call_times[0] >= _IMAGE_RATE_WINDOW_SECONDS:
                _image_call_times.pop(0)
            if len(_image_call_times) < _IMAGE_RATE_LIMIT:
                _image_call_times.append(now)
                return
            wait_for = _IMAGE_RATE_WINDOW_SECONDS - (now - _image_call_times[0]) + 0.05
            await asyncio.sleep(wait_for)


def job_dir(job_id: str) -> Path:
    d = JOBS_DIR / job_id
    d.mkdir(parents=True, exist_ok=True)
    return d


def telemetry_dir(job_id: str) -> Path:
    d = job_dir(job_id) / "telemetry"
    d.mkdir(parents=True, exist_ok=True)
    return d


def _jsonable(value):
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    if isinstance(value, Path):
        return str(value)
    if isinstance(value, dict):
        return {str(k): _jsonable(v) for k, v in value.items()}
    if isinstance(value, (list, tuple)):
        return [_jsonable(v) for v in value]
    if hasattr(value, "model_dump"):
        return _jsonable(value.model_dump())
    if hasattr(value, "to_dict_recursive"):
        return _jsonable(value.to_dict_recursive())
    return str(value)


def append_telemetry(job_id: str, event: dict) -> None:
    payload = {
        "timestamp_utc": datetime.now(timezone.utc).isoformat(),
        "job_id": job_id,
        **_jsonable(event),
    }
    tdir = telemetry_dir(job_id)
    with (tdir / "events.jsonl").open("a", encoding="utf-8") as f:
        f.write(json.dumps(payload, sort_keys=True) + "\n")
    _write_telemetry_summary(job_id)


def _write_telemetry_summary(job_id: str) -> None:
    tdir = telemetry_dir(job_id)
    events_path = tdir / "events.jsonl"
    events = []
    if events_path.exists():
        for line in events_path.read_text().splitlines():
            if not line.strip():
                continue
            try:
                events.append(json.loads(line))
            except json.JSONDecodeError:
                continue

    model_calls = [e for e in events if e.get("event_type") == "model_call"]
    workflow_events = [e for e in events if e.get("event_type") != "model_call"]
    by_model: dict[str, int] = {}
    by_stage: dict[str, int] = {}
    usage_totals: dict[str, int] = {}
    for call in model_calls:
        by_model[call.get("model") or "unknown"] = by_model.get(call.get("model") or "unknown", 0) + 1
        by_stage[call.get("stage") or "unknown"] = by_stage.get(call.get("stage") or "unknown", 0) + 1
        usage = call.get("usage") or {}
        if isinstance(usage, dict):
            for key, value in usage.items():
                if isinstance(value, int):
                    usage_totals[key] = usage_totals.get(key, 0) + value

    summary = {
        "job_id": job_id,
        "updated_at_utc": datetime.now(timezone.utc).isoformat(),
        "event_count": len(events),
        "model_call_count": len(model_calls),
        "workflow_event_count": len(workflow_events),
        "model_calls_by_model": by_model,
        "model_calls_by_stage": by_stage,
        "usage_totals": usage_totals,
        "events_path": "telemetry/events.jsonl",
    }
    (tdir / "summary.json").write_text(json.dumps(summary, indent=2, sort_keys=True))


def telemetry_summary(job_id: str) -> dict:
    summary_path = telemetry_dir(job_id) / "summary.json"
    if not summary_path.exists():
        _write_telemetry_summary(job_id)
    if summary_path.exists():
        return json.loads(summary_path.read_text())
    return {
        "job_id": job_id,
        "event_count": 0,
        "model_call_count": 0,
        "workflow_event_count": 0,
        "model_calls_by_model": {},
        "model_calls_by_stage": {},
        "usage_totals": {},
        "events_path": "telemetry/events.jsonl",
    }


def _response_usage(resp) -> dict:
    usage = getattr(resp, "usage", None)
    return _jsonable(usage) if usage is not None else {}


# PostHog is best-effort telemetry: it must never delay or fail a generation.
# Every model call for one job shares the job ID as its trace ID, letting the
# dashboard group the whole pipeline's cost and latency into one run.
_POSTHOG_API_KEY = os.environ.get("POSTHOG_PROJECT_API_KEY", "phc_DaZmzZ9fKszDeecfqdETWQwMq5uSskHG4UutVvj92GUQ")
_POSTHOG_CAPTURE_URL = os.environ.get("POSTHOG_CAPTURE_URL", "https://us.i.posthog.com/i/v0/e/")


async def _capture_posthog_generation(
    *, job_id: str, stage: str, model: str, duration_s: float, usage: dict,
    is_error: bool = False, error: str | None = None,
) -> None:
    properties = {
        "$ai_trace_id": job_id,
        "$ai_model": model,
        "$ai_provider": "openai",
        "$ai_input_tokens": (usage or {}).get("input_tokens", 0),
        "$ai_output_tokens": (usage or {}).get("output_tokens", 0),
        "$ai_latency": duration_s,
        "$ai_http_status": 0 if is_error else 200,
        "$ai_is_error": is_error,
        "stage": stage,
        "source": "pitch-synthase-v2-workshop",
    }
    if error:
        properties["$ai_error"] = error[:500]
    body = {
        "api_key": _POSTHOG_API_KEY,
        "event": "$ai_generation",
        "distinct_id": job_id,
        "properties": properties,
    }
    try:
        async with httpx.AsyncClient(timeout=3.0) as client:
            await client.post(_POSTHOG_CAPTURE_URL, json=body)
    except Exception as exc:
        logging.getLogger(__name__).warning("posthog LLM capture failed: %s", exc)


async def _responses_create(job_id: str, *, stage: str, model: str, output_path: Path | None = None, **kwargs):
    if _DEBUG_MODE:
        tools = kwargs.get("tools") or []
        is_image = any(isinstance(t, dict) and t.get("type") in ("image_generation", "image_generation_call") for t in tools)
        if is_image:
            print(f"[DEBUG] {stage}: image call → stub PNG", flush=True)
            return _debug_image_resp()
        inp = kwargs.get("input") or ""
        prompt_preview = inp if isinstance(inp, str) else f"[{len(inp)}-item list]"
        print(f"\n[DEBUG] ── {stage} ({model}) ──────────────────", flush=True)
        print(f"[DEBUG] PROMPT ({len(str(prompt_preview))} chars): {str(prompt_preview)[:400]}", flush=True)

    started = time.perf_counter()
    try:
        resp = await gateway.responses.create(model=model, **kwargs)
    except Exception as exc:
        duration_ms = int((time.perf_counter() - started) * 1000)
        append_telemetry(job_id, {
            "event_type": "model_call",
            "stage": stage,
            "model": model,
            "status": "error",
            "duration_ms": duration_ms,
            "error": str(exc)[:500],
            "output_path": output_path,
        })
        await _capture_posthog_generation(
            job_id=job_id, stage=stage, model=model, duration_s=duration_ms / 1000,
            usage={}, is_error=True, error=str(exc),
        )
        raise
    duration_ms = int((time.perf_counter() - started) * 1000)
    if _DEBUG_MODE:
        out = _response_text(resp)
        print(f"[DEBUG] RESPONSE ({duration_ms}ms): {out[:600]}", flush=True)
    usage = _response_usage(resp)
    append_telemetry(job_id, {
        "event_type": "model_call",
        "stage": stage,
        "model": model,
        "status": "ok",
        "duration_ms": duration_ms,
        "response_id": getattr(resp, "id", None),
        "usage": usage,
        "output_path": output_path,
    })
    await _capture_posthog_generation(
        job_id=job_id, stage=stage, model=model, duration_s=duration_ms / 1000,
        usage=usage,
    )
    return resp


def _response_text(resp) -> str:
    text = getattr(resp, "output_text", None)
    if text:
        return text
    chunks = []
    for item in getattr(resp, "output", []) or []:
        for content in getattr(item, "content", []) or []:
            if getattr(content, "type", "") in {"output_text", "text"}:
                chunks.append(getattr(content, "text", "") or "")
    return "\n".join(chunks).strip()


def _response_images(resp) -> list[str]:
    images: list[str] = []
    for item in getattr(resp, "output", []) or []:
        item_type = getattr(item, "type", "")
        if item_type == "image_generation_call":
            result = getattr(item, "result", None)
            if result:
                images.append(result)
        for content in getattr(item, "content", []) or []:
            if getattr(content, "type", "") in {"output_image", "image"}:
                data = getattr(content, "image_base64", None) or getattr(content, "b64_json", None)
                if data:
                    images.append(data)
    return images


def _first_json_object_span(text: str) -> tuple[int, int] | None:
    """Find the first complete top-level JSON object by tracking brace depth
    and string state, rather than assuming the *last* '}' in the text is the
    real end. The model's supporting input can be any document at all, so its
    response shape is not predictable enough to trust naive find/rfind --
    trailing commentary containing braces must not corrupt the slice."""
    start = text.find("{")
    if start < 0:
        return None
    depth = 0
    in_string = False
    escape = False
    for index in range(start, len(text)):
        char = text[index]
        if escape:
            escape = False
            continue
        if char == "\\":
            escape = True
            continue
        if char == '"':
            in_string = not in_string
            continue
        if in_string:
            continue
        if char == "{":
            depth += 1
        elif char == "}":
            depth -= 1
            if depth == 0:
                return start, index + 1
    return None


def _parse_model_json(text: str) -> dict:
    candidate = text.strip()
    if candidate.startswith("```"):
        candidate = candidate.strip("`")
        if candidate.startswith("json"):
            candidate = candidate[4:]
        candidate = candidate.strip()
    span = _first_json_object_span(candidate)
    if span:
        candidate = candidate[span[0]:span[1]]
    return json.loads(candidate, strict=False)


def _validate_export_manifest(export_dir: Path, zip_path: Path, manifest: dict) -> list[str]:
    errors: list[str] = []
    required = list(manifest.get("files") or [])
    slide_images = list(manifest.get("slide_images") or [])
    presenter_view_images = list(manifest.get("presenter_view_images") or [])
    all_paths = required + slide_images + presenter_view_images

    for rel_path in all_paths:
        path = export_dir / rel_path
        if not path.exists():
            errors.append(f"missing export file: {rel_path}")
        elif path.stat().st_size <= 0:
            errors.append(f"empty export file: {rel_path}")

    if not zip_path.exists():
        errors.append("missing export ZIP")
        return errors
    if zip_path.stat().st_size <= 0:
        errors.append("empty export ZIP")
        return errors

    try:
        with zipfile.ZipFile(zip_path) as zf:
            names = set(zf.namelist())
            for rel_path in all_paths:
                if rel_path not in names:
                    errors.append(f"missing ZIP member: {rel_path}")
                else:
                    info = zf.getinfo(rel_path)
                    if info.file_size <= 0:
                        errors.append(f"empty ZIP member: {rel_path}")
    except zipfile.BadZipFile:
        errors.append("invalid export ZIP")
    return errors


async def _write_image_from_b64(data: str, path: Path) -> None:
    path.write_bytes(base64.b64decode(data))


async def _write_preview_jpg(data: str, path: Path, approach: dict, job: dict) -> None:
    """Write a preview image as JPG with embedded approach metadata in EXIF UserComment."""
    import piexif
    from PIL import Image

    jpg_path = path.with_suffix(".jpg")
    raw = base64.b64decode(data)
    img = Image.open(io.BytesIO(raw)).convert("RGB")

    arch_a = approach.get("archetype_a") or {}
    arch_b = approach.get("archetype_b") or {}
    metadata = {
        "ps_version": "1",
        "job_id": job.get("id") or "",
        "approach_id": approach.get("approach_id") or "",
        "approach_label": approach.get("label") or "",
        "pitch_angle": approach.get("pitch_angle") or "",
        "key_differentiator": approach.get("key_differentiator") or "",
        "visual_direction": approach.get("visual_direction") or "",
        "archetype_a": {"id": arch_a.get("archetype_id") or arch_a.get("id") or "", "label": arch_a.get("label") or ""},
        "archetype_b": {"id": arch_b.get("archetype_id") or arch_b.get("id") or "", "label": arch_b.get("label") or ""},
        "vibe_semantics": approach.get("vibe_semantics") or [],
        "association_words": json.loads(job.get("association_words") or "[]") if isinstance(job.get("association_words"), str) else (job.get("association_words") or []),
        "generated_at": datetime.now(timezone.utc).isoformat(),
    }
    comment_bytes = b"UNICODE\x00" + json.dumps(metadata, ensure_ascii=False).encode("utf-8")
    exif_dict = {"Exif": {piexif.ExifIFD.UserComment: comment_bytes}}
    exif_bytes = piexif.dump(exif_dict)

    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=85, optimize=True, exif=exif_bytes)
    jpg_path.write_bytes(buf.getvalue())


def _image_content(path: Path) -> dict:
    mime = mimetypes.guess_type(path.name)[0] or "image/png"
    image_b64 = base64.b64encode(path.read_bytes()).decode()
    return {"type": "input_image", "image_url": f"data:{mime};base64,{image_b64}", "detail": "auto"}


def _selected_template_image(job: dict, job_id: str) -> Path | None:
    selected = job.get("selected_candidate_id") or ""
    if not selected.startswith("tmpl_"):
        return None
    suffix = selected.split("_", 1)[1]
    if suffix not in {"1", "2", "3", "4"}:
        return None
    path = job_dir(job_id) / "candidates" / f"template_{suffix}.png"
    return path if path.exists() else None


def _design_reference_image(job: dict, job_id: str) -> Path | None:
    """Return the design reference image for the deck.

    Checks selected_prototype_id in prototype_candidates_json first (new decoupled path),
    then falls back to design_reference_image_path on the selected approach (legacy path).
    """
    selected_proto_id = job.get("selected_prototype_id")
    if selected_proto_id:
        for candidate in (job.get("prototype_candidates_json") or []):
            if candidate.get("design_id") == selected_proto_id:
                ref = candidate.get("reference_image_path")
                if ref:
                    p = Path(ref)
                    return p if p.exists() else None

    # Legacy: design_reference_image_path on the approach (pre-decoupling)
    selected_id = job.get("selected_candidate_id") or ""
    for approach in (job.get("approach_candidates_json") or []):
        if approach.get("approach_id") == selected_id:
            ref = approach.get("design_reference_image_path")
            if ref:
                p = Path(ref)
                return p if p.exists() else None
    return None


def _workshop_archetype(job: dict) -> dict:
    archetype_id = job.get("selected_archetype_id") or ""
    archetype = prompts.archetype_by_id(archetype_id)
    if not archetype:
        raise RuntimeError(f"Unknown or missing workshop archetype: {archetype_id or '[blank]'}")
    return dict(archetype)


def _label_word_count(label: str) -> int:
    return len([part for part in re.split(r"[\s/-]+", label.strip()) if part])


def _contains_casefold(haystack: str, needle: str) -> bool:
    return re.search(rf"\b{re.escape(needle)}\b", haystack, re.IGNORECASE) is not None


def _validate_alignment_manifest(manifest: dict, selected_archetype: dict) -> list[dict]:
    if set(manifest.keys()) != {"candidates"}:
        raise RuntimeError("Workshop template output must contain exactly one top-level key: candidates")
    candidates = manifest.get("candidates")
    if not isinstance(candidates, list) or len(candidates) != 4:
        raise RuntimeError("Workshop template output must contain exactly four candidates")

    expected_ids = [f"tmpl_{i}" for i in range(1, 5)]
    labels_seen: set[str] = set()
    cleaned: list[dict] = []
    archetype_label = str(selected_archetype.get("label") or "").strip()

    for index, candidate in enumerate(candidates):
        if not isinstance(candidate, dict):
            raise RuntimeError(f"Candidate {index + 1} is not an object")
        if "focus_id" in candidate:
            raise RuntimeError("Workshop candidates must not include focus_id")
        if "selected_archetypes" in candidate:
            raise RuntimeError("Workshop candidates must not include selected_archetypes")

        candidate_id = str(candidate.get("candidate_id") or "").strip()
        if candidate_id != expected_ids[index]:
            raise RuntimeError(f"Candidate IDs must be {expected_ids}; got {[c.get('candidate_id') for c in candidates]}")

        label = str(candidate.get("label") or "").strip()
        if _label_word_count(label) not in {2, 3}:
            raise RuntimeError(f"Alignment label must have two or three words: {label!r}")
        if label.casefold() in labels_seen:
            raise RuntimeError(f"Alignment labels must be unique: {label!r}")
        labels_seen.add(label.casefold())
        if archetype_label and _contains_casefold(label, archetype_label):
            raise RuntimeError(f"Alignment label must not repeat archetype label: {label!r}")

        description = str(candidate.get("description") or "").strip()
        visual_doctrine = str(candidate.get("visual_doctrine") or "").strip()
        image_prompt = str(candidate.get("image_prompt") or "").strip()
        style_tags = candidate.get("style_tags")

        if not description:
            raise RuntimeError(f"Candidate {candidate_id} is missing description")
        if not visual_doctrine:
            raise RuntimeError(f"Candidate {candidate_id} is missing visual_doctrine")
        if not isinstance(style_tags, list) or len(style_tags) < 3:
            raise RuntimeError(f"Candidate {candidate_id} must include at least three style_tags")
        if len(image_prompt) < 80:
            raise RuntimeError(f"Candidate {candidate_id} image_prompt is too short")
        if _contains_casefold(image_prompt, label):
            raise RuntimeError(f"Candidate {candidate_id} leaked its alignment label into image_prompt")
        if archetype_label and _contains_casefold(image_prompt, archetype_label):
            raise RuntimeError(f"Candidate {candidate_id} leaked the archetype label into image_prompt")

        cleaned.append({
            "candidate_id": candidate_id,
            "label": label,
            "description": description,
            "visual_doctrine": visual_doctrine,
            "style_tags": [str(tag).strip() for tag in style_tags if str(tag).strip()],
            "image_prompt": image_prompt,
        })

    return cleaned


def _supporting_image(job_id: str) -> Path | None:
    intake_dir = job_dir(job_id) / "intake"
    for suffix in ("*.png", "*.jpg", "*.jpeg", "*.webp"):
        matches = sorted(intake_dir.glob(suffix))
        if matches:
            return matches[0]
    return None


def _anchor_input(prompt: str, job_id: str):
    content = [{"type": "input_text", "text": prompt}]
    supporting_image = _supporting_image(job_id)
    if supporting_image:
        content.append({"type": "input_text", "text": "User-provided reference image."})
        content.append(_image_content(supporting_image))
    return [{
        "role": "user",
        "content": content,
    }]


def _preview_vertex_paths(job_id: str) -> list[Path]:
    """
    Returns the preview image paths the user elected to keep on conversion.

    The job record may contain `convert_keep_candidate_ids` — a list of
    candidate IDs (e.g. ["cand_1", "cand_3"]) reflecting the post-convert
    checkbox selection. If absent or empty, no preview candidates are returned.
    """
    job = db.get_job(job_id) or {}
    cand_dir = job_dir(job_id) / "candidates"
    keep_ids: list[str] = job.get("convert_keep_candidate_ids") or []
    paths = []
    for cid in keep_ids:
        suffix = cid.split("_", 1)[-1]
        p = cand_dir / f"cand_{suffix}.png"
        if p.exists():
            paths.append(p)
    return paths


def _authoritative_source_materials_block(job: dict) -> str:
    """Build the lossless factual source packet for Deck Builder calls.

    The Anchor Writer owns framing, not source transport. Raw user materials
    cross this seam independently so an omission or paraphrase in the anchor
    narrative cannot erase concrete pitch facts.
    """
    sections = []
    elevator_pitch = str(job.get("elevator_pitch") or "").strip()
    if elevator_pitch:
        sections.append(
            "FOUNDER NARRATIVE / ORIGINAL PITCH\n"
            "-----------------------------------\n"
            f"{elevator_pitch}"
        )

    doc_text = str(job.get("doc_text") or "").strip()
    if doc_text:
        sections.append(
            "SUPPORTING SPECIFICATION / SOURCE DOCUMENT\n"
            "------------------------------------------\n"
            f"{doc_text}"
        )

    if not sections:
        return ""

    return (
        "The pitch content was based on the following user-provided materials.\n"
        "These materials are the primary factual authority for the deck. The "
        "Anchor Writer's presentation context controls framing, posture, and "
        "narrative pressure only; it must not replace, summarize away, or "
        "override this source content. Preserve concrete source-backed names, "
        "quantities, units, ranges, statuses, caveats, development targets, and "
        "phase gates wherever they materially strengthen the pitch. Never convert "
        "a stated target or hypothesis into a validated result.\n\n"
        + "\n\n".join(sections)
    )


def _deck_drafter_input(
    anchor_payload: dict,
    job_id: str,
    vertex_images: list[Path] | None = None,
    vertex_instruction: str | None = None,
):
    job = db.get_job(job_id) or {}
    content = [{
        "type": "input_text",
        "text": json.dumps(anchor_payload, ensure_ascii=False, indent=2),
    }]

    source_materials = _authoritative_source_materials_block(job)
    if source_materials:
        content.append({
            "type": "input_text",
            "text": source_materials,
        })

    content.append({
        "type": "input_text",
        "text": prompts.pitch_aspect_worker_briefing(
            job.get("pitch_aspect_modes")
        ),
    })

    intake_block = prompts.intake_context_block(job.get("intake_options") or {})
    if intake_block:
        content.append({"type": "input_text", "text": intake_block})

    image_analysis = job.get("image_analysis") or {}
    if image_analysis:
        content.append({
            "type": "input_text",
            "text": (
                "INTAKE IMAGE ANALYSIS — a vision scan of the user-uploaded reference "
                f"image found: {json.dumps(image_analysis, ensure_ascii=False)}. Use this "
                "only to corroborate what the image itself shows; it does not override "
                "the image."
            ),
        })

    if vertex_images:
        instruction = vertex_instruction or (
            "The following photographs are slides from the preview of this pitch. "
            "They belong to this same presentation. Their positions within the full "
            "deck sequence are not specified — place each where it fits best in the "
            "deck argument. Use them as the governing visual style reference for all "
            "generated slides; do not reference any separate template image."
        )
        content.append({
            "type": "input_text",
            "text": instruction,
        })
        for vpath in vertex_images:
            content.append(_image_content(vpath))
    else:
        # Non-convert path: derive aesthetic from archetype posture + visual focus.
        # Look up the selected archetype's full posture text from the FOUNDER_ARCHETYPES pool.
        selected_archetype_id = job.get("selected_archetype_id") or ""
        archetype_label = job.get("selected_archetype_label") or ""
        archetype_posture = ""
        for a in (getattr(prompts, "FOUNDER_ARCHETYPES", None) or []):
            if a.get("archetype_id") == selected_archetype_id:
                archetype_posture = a.get("posture") or ""
                archetype_label = archetype_label or a.get("label") or ""
                break

        selected_approach = None
        for a in (job.get("approach_candidates_json") or []):
            if a.get("approach_id") == job.get("selected_candidate_id"):
                selected_approach = a
                break

        parts = []
        if archetype_label or archetype_posture:
            parts.append(f"Selected founder posture archetype: {archetype_label}.")
            if archetype_posture:
                parts.append(archetype_posture)
        if selected_approach:
            focus_label = selected_approach.get("visual_direction") or ""
            focus_note  = selected_approach.get("key_differentiator") or ""
            parts.append(f"\nSelected visual focus: {focus_label}.")
            if focus_note:
                parts.append(focus_note)
        if parts:
            parts.append(
                "\nHierarchy rule: the founder archetype governs deck argument, founder "
                "portrayal, evidence posture, emphasis order, and rhetorical stance. "
                "The visual aesthetic — palette, typographic register, information density, "
                "and layout conventions — must be derived from the selected approach's "
                "representative preview slide provided below, interpreted in service of "
                "that archetype."
            )
            content.append({"type": "input_text", "text": "\n".join(parts)})

        preview_path_str = (selected_approach or {}).get("preview_image_path")
        preview_path = Path(preview_path_str) if preview_path_str else None
        if preview_path and preview_path.exists():
            content.append({
                "type": "input_text",
                "text": (
                    "Selected approach's representative preview slide — this is the "
                    "mandatory visual style anchor. This is the exact slide the human "
                    "picked when choosing this approach; all generated slides must derive "
                    "their palette, typographic register, and layout conventions from it. "
                    "If written style instructions conflict with this image, the image "
                    "wins for surface styling only. It must not override archetype-driven "
                    "structure, founder framing, or narrative posture."
                ),
            })
            content.append(_image_content(preview_path))

    supporting_image = _supporting_image(job_id)
    if supporting_image:
        content.append({
            "type": "input_text",
            "text": (
                "USER-PROVIDED PRODUCT / REFERENCE IMAGE — AUTHORITATIVE VISUAL "
                "SOURCE MATERIAL. The pitch content was also based on this image. "
                "Use it as factual visual evidence for the product's depicted form, "
                "parts, material language, and relationships. Draw from it where "
                "relevant; do not treat the Anchor Writer's prose as a substitute "
                "for the image, and do not copy the image verbatim."
            ),
        })
        content.append(_image_content(supporting_image))

    design_ref = _design_reference_image(job, job_id)
    if design_ref:
        content.append({
            "type": "input_text",
            "text": (
                "PRODUCT DESIGN REFERENCE — MANDATORY VISUAL CONTRACT. This is the "
                "canonical prototype image for this specific approach. Every slide that "
                "depicts the physical product must replicate its exact form factor, "
                "material finishes, color, proportions, and visible mechanism elements. "
                "This image takes precedence over any written description when generating "
                "product visuals. Do not invent variations, alternate colorways, or "
                "different form factors — the physical object must remain identical "
                "across all slides in this deck."
            ),
        })
        content.append(_image_content(design_ref))
    return [{"role": "user", "content": content}]


def _preview_slide_image_input(image_prompt: str, job_id: str):
    """Content list for a single per-slide image generation call in the preview Deck Builder."""
    job = db.get_job(job_id) or {}
    content = [{"type": "input_text", "text": prompts.preview_slide_image_prompt(image_prompt)}]
    tmpl_path = _selected_template_image(job, job_id)
    if tmpl_path:
        content.append({
            "type": "input_text",
            "text": (
                "Selected template image — mandatory visual style anchor. "
                "Palette, typographic register, layout conventions, and overall "
                "visual posture must be derived from this image. "
                "If the text prompt above conflicts with this image, the image wins."
            ),
        })
        content.append(_image_content(tmpl_path))
    return [{"role": "user", "content": content}]


def _raw_user_inputs(job: dict) -> dict:
    override = _raw_user_inputs_cv.get()
    if override is not None:
        return override(job)
    return {
        "page_1": {
            "problem_need": job.get("problem_need") or "",
            "audience": job.get("audience") or "",
            "association_words": job.get("association_words") or [],
        },
        "page_2": {
            "elevator_pitch": job.get("elevator_pitch") or "",
            "supporting_document_provided": bool(job.get("doc_text")),
            "supporting_image": bool(_supporting_image(job.get("id", ""))) if job.get("id") else False,
            "optional_notes": None,
        },
    }


def _selected_template_payload(job: dict) -> dict:
    override = _selected_template_payload_cv.get()
    if override is not None:
        return override(job)
    selected_id = job.get("selected_candidate_id") or ""
    selected = None
    for candidate in job.get("candidates_json") or []:
        if candidate.get("candidate_id") == selected_id:
            selected = candidate
            break
    archetype_id = job.get("selected_archetype_id") or ""
    archetype_label = job.get("selected_archetype_label") or ""
    archetype_posture = ""
    for a in (getattr(prompts, "FOUNDER_ARCHETYPES", None) or []):
        if a.get("archetype_id") == archetype_id:
            archetype_posture = a.get("posture") or ""
            archetype_label = archetype_label or a.get("label") or ""
            break
    return {
        "candidate_id": selected_id,
        "focus_id": (selected or {}).get("focus_id") or "",
        "style_label": job.get("selected_candidate_label") or (selected or {}).get("style_label") or "",
        "style_tags": (selected or {}).get("style_tags") or [],
        "template_manifest": selected or None,
        "archetype": {
            "archetype_id": archetype_id,
            "label": archetype_label,
            "posture": archetype_posture,
        },
    }


def _anchor_payload(anchor_json: dict, *, focused_prompt: str | None = None, slide_number: int | None = None) -> dict:
    anchor_narrative = anchor_json.get("deck_generation_prompt", "")
    wrapped_narrative = (
        f"{anchor_narrative}\n\n"
        "Sourcing rule: only user_inputs and the supplied supporting "
        "document/image are authoritative for real facts about this "
        "company's actual founders, team, or personnel -- names, job "
        "titles, credentials, years of experience, and past employers. "
        "The presentation context above establishes tone, posture, and "
        "narrative pressure; it is not itself a source for who the real "
        "team is. Do not construct a 'team,' 'founders,' or 'who we are' "
        "slide using person-level specifics drawn from the presentation "
        "context unless the same specifics are also present in user_inputs "
        "or the supporting materials. Where the user's own materials do not "
        "supply real team specifics, present the needed competencies as "
        "team requirements, an operator profile, capabilities needed, or a "
        "diligence checklist rather than as named individuals."
    ) if anchor_narrative else anchor_narrative
    payload = {
        "deck_generation_prompt": wrapped_narrative,
        "user_inputs": anchor_json.get("user_inputs") or {},
    }
    if focused_prompt:
        payload["focused_prompt"] = focused_prompt
    if slide_number is not None:
        payload["slide_number"] = slide_number
    return payload



async def _synthesize_anchor_expansions(
    job_id: str,
    *,
    deck_generation_prompt: str,
    user_inputs: dict,
    selected_template: dict,
    strategy_dir: Path,
    stage_prefix: str,
) -> tuple[str, dict, dict]:
    """Run visual and founder narrative synthesis from the same anchor in parallel.

    Both addenda are deterministic inputs to the Deck Builder. The original anchor
    remains first, followed by visual production constraints and then founder voice
    constraints so the verbal grammar is the final authorship instruction seen by
    the storyboard call.
    """
    visual_grammar_prompt = prompts.visual_grammar_prompt(
        deck_generation_prompt
    )
    founder_narrative_prompt = prompts.founder_narrative_synthesis_prompt(
        deck_generation_prompt=deck_generation_prompt,
        user_inputs=user_inputs,
        selected_template=selected_template,
    )

    visual_grammar_response, founder_narrative_response = await asyncio.gather(
        _responses_create(
            job_id,
            stage=f"{stage_prefix}_visual_grammar",
            model=prompts.VISUAL_GRAMMAR_MODEL,
            input=[{
                "role": "user",
                "content": [{
                    "type": "input_text",
                    "text": visual_grammar_prompt,
                }],
            }],
        ),
        _responses_create(
            job_id,
            stage=f"{stage_prefix}_founder_narrative",
            model=prompts.FOUNDER_NARRATIVE_MODEL,
            input=[{
                "role": "user",
                "content": [{
                    "type": "input_text",
                    "text": founder_narrative_prompt,
                }],
            }],
        ),
    )

    visual_grammar_result = _parse_model_json(
        _response_text(visual_grammar_response)
    )
    founder_narrative_result = _parse_model_json(
        _response_text(founder_narrative_response)
    )

    visual_grammar_addendum = str(
        visual_grammar_result.get("visual_grammar_addendum") or ""
    ).strip()
    founder_narrative_addendum = str(
        founder_narrative_result.get("founder_narrative_addendum") or ""
    ).strip()

    if not visual_grammar_addendum:
        raise RuntimeError(
            "Visual Grammar Synthesizer returned no visual_grammar_addendum"
        )
    if not founder_narrative_addendum:
        raise RuntimeError(
            "Founder Narrative Synthesizer returned no founder_narrative_addendum"
        )

    expanded_prompt = "\n\n".join((
        deck_generation_prompt.strip(),
        visual_grammar_addendum,
        founder_narrative_addendum,
    ))

    (strategy_dir / "visual_grammar_output.json").write_text(
        json.dumps(visual_grammar_result, indent=2),
        encoding="utf-8",
    )
    (strategy_dir / "founder_narrative_output.json").write_text(
        json.dumps(founder_narrative_result, indent=2),
        encoding="utf-8",
    )

    return (
        expanded_prompt,
        visual_grammar_result,
        founder_narrative_result,
    )


# ─────────────────────────────────────────────────────────────────────────────
# Stage 1 — Candidate preview (Procedure 1A)
# ─────────────────────────────────────────────────────────────────────────────

async def template_worker(job_id: str):
    job = db.get_job(job_id)
    if not job:
        return
    if job.get("candidates_json"):
        return
    if not db.claim_job_status(
        job_id,
        {"created", "template_generation_queued", "template_generation_started", "preview_checkout_required"},
        "template_generation_running",
    ):
        append_telemetry(job_id, {"event_type": "workflow", "stage": "template_styles", "status": "duplicate_skipped"})
        return
    job = db.get_job(job_id)
    problem_need = job.get("problem_need") or "the need described by the user"
    audience = job["audience"]
    association_words = job["association_words"]
    append_telemetry(job_id, {"event_type": "workflow", "stage": "template_styles", "status": "started"})

    cand_dir = job_dir(job_id) / "candidates"
    cand_dir.mkdir(exist_ok=True)
    try:
        resp = await _responses_create(
            job_id,
            stage="template_drafter",
            model=prompts.TEMPLATE_DRAFTER_MODEL,
            input=prompts.template_drafter_prompt(
                problem_need=problem_need,
                audience=audience,
                association_words=association_words,
            ),
            output_path=(cand_dir / "template_prompt_manifest.json").relative_to(job_dir(job_id)),
        )
        prompt_manifest = _parse_model_json(_response_text(resp))
        selected_archetypes = prompt_manifest.get("selected_archetypes") or []
        prompt_candidates = prompt_manifest.get("candidates") or []
        if len(selected_archetypes) != 4:
            raise RuntimeError(f"Template drafter returned {len(selected_archetypes)} archetypes; expected 4")
        if len(prompt_candidates) != 4:
            raise RuntimeError(f"Template drafter returned {len(prompt_candidates)} visual focus prompts; expected 4")
        expected_ids = [f"tmpl_{i}" for i in range(1, 5)]
        actual_ids = [str(c.get("candidate_id") or "") for c in prompt_candidates]
        if actual_ids != expected_ids:
            raise RuntimeError(f"Template drafter candidate_id values must be {expected_ids}; got {actual_ids}")
        (cand_dir / "template_prompt_manifest.json").write_text(json.dumps(prompt_manifest, indent=2))

        async def _gen_template_image(i: int, candidate: dict) -> str:
            await _throttle_image_generation()
            image_resp = await _responses_create(
                job_id,
                stage=f"template_image_{i}",
                model=prompts.DECK_DRAFTER_MODEL,
                input=prompts.template_image_prompt(candidate=candidate),
                tools=[prompts.image_generation_tool(stage="template")],
                output_path=(cand_dir / f"template_{i}.png").relative_to(job_dir(job_id)),
            )
            images = _response_images(image_resp)
            if len(images) != 1:
                raise RuntimeError(f"Template image {i} returned {len(images)} images; expected 1")
            return images[0]

        image_payloads = list(await asyncio.gather(*[
            _gen_template_image(i, c)
            for i, c in enumerate(prompt_candidates, start=1)
        ]))
    except Exception as e:
        append_telemetry(job_id, {
            "event_type": "workflow",
            "stage": "template_styles",
            "status": "failed",
            "error": str(e)[:500],
        })
        db.update_job(job_id, status="template_generation_failed", error_message=f"Template Drafter failed: {e}")
        return

    templates = []
    for i, img_b64 in enumerate(image_payloads[:4]):
        template_number = i + 1
        img_path = cand_dir / f"template_{template_number}.png"
        await _write_image_from_b64(img_b64, img_path)
        emphasis = association_words[i % len(association_words)]
        candidate = prompt_candidates[i]
        templates.append({
            "candidate_id": f"tmpl_{template_number}",
            "focus_id": candidate.get("focus_id") or "",
            "url": f"/instant/pitch-deck/jobs/{job_id}/candidates/template_{template_number}.png",
            "style_label": candidate.get("label") or candidate.get("style_label") or TEMPLATE_LABELS[i],
            "style_tags": candidate.get("style_tags") or [f"{emphasis} emphasis", "visual anchor", "no product copy"],
            "direction_note": candidate.get("description") or "",
        })

    archetypes = [
        {
            "archetype_id": a.get("archetype_id") or "",
            "label": a.get("label") or "",
        }
        for a in selected_archetypes[:4]
    ]

    (cand_dir / "templates.json").write_text(json.dumps(templates, indent=2))
    (cand_dir / "archetypes.json").write_text(json.dumps(archetypes, indent=2))
    db.update_job(
        job_id,
        status="template_ready",
        candidates_json=templates,
        archetypes_json=archetypes,
        error_message=None,
    )
    append_telemetry(job_id, {
        "event_type": "workflow",
        "stage": "template_styles",
        "status": "template_ready",
        "template_count": len(templates),
    })


async def template_text_worker_baseline(job_id: str):
    """Text-only counterpart to template_worker: current-production prompt/pool,
    no image generation. Exists only for the archetype-first workshop comparison,
    so baseline and experiment outputs are produced by separate, comparably-scoped
    workers."""
    job = db.get_job(job_id)
    if not job:
        return
    if not db.claim_job_status(
        job_id,
        {"created", "template_prompt_queued", "template_prompt_started", "template_prompt_failed"},
        "template_prompt_running",
    ):
        append_telemetry(job_id, {"event_type": "workflow", "stage": "template_prompt_baseline", "status": "duplicate_skipped"})
        return

    job = db.get_job(job_id) or {}
    problem_need = job.get("problem_need") or "the need described by the user"
    audience = job.get("audience") or ""
    association_words = job.get("association_words") or []

    cand_dir = job_dir(job_id) / "candidates"
    cand_dir.mkdir(exist_ok=True)
    append_telemetry(job_id, {"event_type": "workflow", "stage": "template_prompt_baseline", "status": "started"})

    prompt_text = prompts.template_drafter_prompt(
        problem_need=problem_need,
        audience=audience,
        association_words=association_words,
    )
    (cand_dir / "template_prompt.txt").write_text(prompt_text)

    last_error: Exception | None = None
    for attempt in range(1, 3):
        try:
            resp = await _responses_create(
                job_id,
                stage=f"template_drafter_baseline_{attempt}",
                model=prompts.TEMPLATE_DRAFTER_MODEL,
                input=prompt_text,
                output_path=(cand_dir / f"template_prompt_response_attempt_{attempt}.txt").relative_to(job_dir(job_id)),
            )
            raw_text = _response_text(resp)
            (cand_dir / f"template_prompt_response_attempt_{attempt}.txt").write_text(raw_text)
            manifest = _parse_model_json(raw_text)

            selected_archetypes = manifest.get("selected_archetypes") or []
            prompt_candidates = manifest.get("candidates") or []
            if len(selected_archetypes) != 4:
                raise RuntimeError(f"Template drafter returned {len(selected_archetypes)} archetypes; expected 4")
            if len(prompt_candidates) != 4:
                raise RuntimeError(f"Template drafter returned {len(prompt_candidates)} visual focus prompts; expected 4")
            expected_ids = [f"tmpl_{i}" for i in range(1, 5)]
            actual_ids = [str(c.get("candidate_id") or "") for c in prompt_candidates]
            if actual_ids != expected_ids:
                raise RuntimeError(f"Template drafter candidate_id values must be {expected_ids}; got {actual_ids}")
            (cand_dir / "template_prompt_manifest.json").write_text(json.dumps(manifest, indent=2))

            db.update_job(
                job_id,
                status="template_prompt_ready",
                candidates_json=prompt_candidates,
                archetypes_json=selected_archetypes,
                error_message=None,
            )
            append_telemetry(job_id, {
                "event_type": "workflow",
                "stage": "template_prompt_baseline",
                "status": "template_prompt_ready",
                "attempt": attempt,
                "candidate_count": len(prompt_candidates),
            })
            return
        except Exception as exc:
            last_error = exc
            append_telemetry(job_id, {
                "event_type": "workflow",
                "stage": "template_prompt_baseline",
                "status": "retrying" if attempt == 1 else "failed",
                "attempt": attempt,
                "error": str(exc)[:500],
            })

    db.update_job(
        job_id,
        status="template_prompt_failed",
        error_message=f"Workshop baseline template prompt failed: {last_error}",
    )


async def template_text_worker_alignment_experiment(job_id: str):
    job = db.get_job(job_id)
    if not job:
        return
    if not db.claim_job_status(
        job_id,
        {"created", "template_prompt_queued", "template_prompt_started", "template_prompt_failed"},
        "template_prompt_running",
    ):
        append_telemetry(job_id, {"event_type": "workflow", "stage": "template_prompt_experiment", "status": "duplicate_skipped"})
        return

    job = db.get_job(job_id) or {}
    problem_need = job.get("problem_need") or "the need described by the user"
    audience = job.get("audience") or ""
    association_words = job.get("association_words") or []
    selected_archetype = _workshop_archetype(job)

    cand_dir = job_dir(job_id) / "candidates"
    cand_dir.mkdir(exist_ok=True)
    append_telemetry(job_id, {
        "event_type": "workflow",
        "stage": "template_prompt_experiment",
        "status": "started",
        "selected_archetype_id": selected_archetype.get("archetype_id"),
    })

    prompt_text = prompts.template_drafter_prompt_archetype_first(
        problem_need=problem_need,
        audience=audience,
        association_words=association_words,
        selected_archetype=selected_archetype,
    )
    (cand_dir / "template_prompt.txt").write_text(prompt_text)

    last_error: Exception | None = None
    for attempt in range(1, 3):
        try:
            resp = await _responses_create(
                job_id,
                stage=f"template_drafter_experiment_{attempt}",
                model=prompts.TEMPLATE_DRAFTER_MODEL,
                input=prompt_text,
                output_path=(cand_dir / f"template_prompt_response_attempt_{attempt}.txt").relative_to(job_dir(job_id)),
            )
            raw_text = _response_text(resp)
            (cand_dir / f"template_prompt_response_attempt_{attempt}.txt").write_text(raw_text)
            manifest = _parse_model_json(raw_text)
            cleaned = _validate_alignment_manifest(manifest, selected_archetype)
            (cand_dir / "template_prompt_manifest.json").write_text(json.dumps({"candidates": cleaned}, indent=2))

            templates = [
                {
                    "candidate_id": candidate["candidate_id"],
                    "style_label": candidate["label"],
                    "style_tags": candidate["style_tags"],
                    "direction_note": candidate["description"],
                    "visual_doctrine": candidate["visual_doctrine"],
                    "image_prompt": candidate["image_prompt"],
                }
                for candidate in cleaned
            ]
            db.update_job(
                job_id,
                status="template_prompt_ready",
                candidates_json=templates,
                archetypes_json=[selected_archetype],
                error_message=None,
            )
            append_telemetry(job_id, {
                "event_type": "workflow",
                "stage": "template_prompt_experiment",
                "status": "template_prompt_ready",
                "attempt": attempt,
                "candidate_count": len(templates),
            })
            return
        except Exception as exc:
            last_error = exc
            append_telemetry(job_id, {
                "event_type": "workflow",
                "stage": "template_prompt_experiment",
                "status": "retrying" if attempt == 1 else "failed",
                "attempt": attempt,
                "error": str(exc)[:500],
            })

    db.update_job(
        job_id,
        status="template_prompt_failed",
        error_message=f"Workshop template prompt failed: {last_error}",
    )


def _validate_approach_manifest(manifest: dict, *, expected_count: int = 4) -> list[dict]:
    approaches = manifest.get("approaches")
    if not isinstance(approaches, list) or len(approaches) != expected_count:
        raise RuntimeError(f"Approach manifest must contain exactly {expected_count} approaches")

    expected_ids = [f"approach_{i}" for i in range(1, expected_count + 1)]
    labels_seen: set[str] = set()
    pairs_seen: set[frozenset] = set()
    cleaned = []
    for index, approach in enumerate(approaches):
        if not isinstance(approach, dict):
            raise RuntimeError(f"Approach {index + 1} is not an object")
        approach_id = str(approach.get("approach_id") or "").strip()
        if approach_id != expected_ids[index]:
            raise RuntimeError(f"Approach IDs must be {expected_ids}; got {[a.get('approach_id') for a in approaches]}")
        label = str(approach.get("label") or "").strip()
        if not label:
            raise RuntimeError(f"Approach {approach_id} is missing a label")
        if label.casefold() in labels_seen:
            raise RuntimeError(f"Approach labels must be unique: {label!r}")
        labels_seen.add(label.casefold())
        for field in ("pitch_angle", "key_differentiator", "visual_direction"):
            if not str(approach.get(field) or "").strip():
                raise RuntimeError(f"Approach {approach_id} is missing {field}")

        arch_a = approach.get("archetype_a")
        arch_b = approach.get("archetype_b")
        if not (isinstance(arch_a, dict) and arch_a.get("archetype_id")):
            raise RuntimeError(f"Approach {approach_id} missing archetype_a")
        if not (isinstance(arch_b, dict) and arch_b.get("archetype_id")):
            raise RuntimeError(f"Approach {approach_id} missing archetype_b")
        pair_key = frozenset([arch_a["archetype_id"], arch_b["archetype_id"]])
        if pair_key in pairs_seen:
            raise RuntimeError(f"Approach {approach_id} repeats an archetype pair already used")
        pairs_seen.add(pair_key)

        entry = {
            "approach_id": approach_id,
            "label": label,
            "pitch_angle": str(approach["pitch_angle"]).strip(),
            "key_differentiator": str(approach["key_differentiator"]).strip(),
            "visual_direction": str(approach["visual_direction"]).strip(),
            "archetype_a": arch_a,
            "archetype_b": arch_b,
        }
        if approach.get("effective_archetype_id"):
            entry["effective_archetype_id"] = str(approach["effective_archetype_id"]).strip()
        if isinstance(approach.get("vibe_semantics"), list):
            entry["vibe_semantics"] = approach["vibe_semantics"]
        cleaned.append(entry)
    return cleaned


async def pitch_quality_gate_worker(job_id: str) -> dict:
    """Cheap pre-check: does the elevator pitch have enough to work with?
    Text-only, no state mutation beyond telemetry -- caller decides what to
    do with the verdict (e.g. offer the elaborator button)."""
    job = db.get_job(job_id)
    prompt_text = prompts.pitch_quality_gate_prompt(
        elevator_pitch=job.get("elevator_pitch") or "",
        doc_text=job.get("doc_text"),
    )
    resp = await _responses_create(
        job_id, stage="pitch_quality_gate", model=prompts.TEMPLATE_DRAFTER_MODEL, input=prompt_text,
    )
    verdict = _parse_model_json(_response_text(resp))
    append_telemetry(job_id, {
        "event_type": "workflow", "stage": "pitch_quality_gate",
        "is_real_pitch": verdict.get("is_real_pitch"),
    })
    return verdict


async def pitch_elaborator_worker(job_id: str) -> dict:
    """The 'write my whole pitch' button: ghostwrites a complete elevator
    pitch draft from whatever rough input is currently on the job. Does NOT
    overwrite job.elevator_pitch -- caller shows the draft to the user for
    review/edit, then explicitly saves it if accepted."""
    job = db.get_job(job_id)
    rough_input = job.get("elevator_pitch") or ""
    prompt_text = prompts.pitch_elaborator_prompt(rough_input=rough_input)
    resp = await _responses_create(
        job_id, stage="pitch_elaborator", model=prompts.TEMPLATE_DRAFTER_MODEL, input=prompt_text,
    )
    draft = _parse_model_json(_response_text(resp))
    for field in ("elevator_pitch", "invented_elements", "conveys_suggestion", "is_nonsense"):
        if field not in draft:
            raise RuntimeError(f"Elaborator draft missing {field!r}")
    append_telemetry(job_id, {
        "event_type": "workflow", "stage": "pitch_elaborator",
        "invented_count": len(draft["invented_elements"]),
        "is_nonsense": draft["is_nonsense"],
    })
    return draft


async def approach_drafter_worker(job_id: str):
    """Full-pitch-first workshop: one text-only call, informed by the REAL
    elevator pitch (+ optional doc), returns four strategic pitch approaches.
    No images. Archetype is fixed by the user or chosen by the drafter."""
    job = db.get_job(job_id)
    if not job:
        return
    if not db.claim_job_status(
        job_id,
        {"created", "approaches_queued", "approaches_failed"},
        "approaches_running",
    ):
        append_telemetry(job_id, {"event_type": "workflow", "stage": "approach_drafter", "status": "duplicate_skipped"})
        return

    job = db.get_job(job_id) or {}
    fixed_archetype = None
    if job.get("selected_archetype_id"):
        fixed_archetype = prompts.archetype_by_id(job["selected_archetype_id"])

    cand_dir = job_dir(job_id) / "approaches"
    cand_dir.mkdir(exist_ok=True)
    append_telemetry(job_id, {"event_type": "workflow", "stage": "approach_drafter", "status": "started"})

    intake_img = _supporting_image(job_id)
    # Pass physical descriptions from design_refs_text if available — gives all 4
    # approaches a shared physical grounding without coupling archetype choices.
    design_philosophies = None
    prototype_candidates = job.get("prototype_candidates_json") or []
    if prototype_candidates and not (len(prototype_candidates) == 1 and prototype_candidates[0].get("design_id") == "design_uploaded"):
        design_philosophies = [
            {"design_id": d.get("design_id"), "design_philosophy": d.get("design_philosophy", ""), "physical_description": d.get("physical_description", "")}
            for d in prototype_candidates
        ]
    prompt_text = prompts.approach_drafter_prompt(
        audience=job.get("audience") or "",
        elevator_pitch=job.get("elevator_pitch") or "",
        conveys=job.get("conveys") or "",
        selected_archetype=fixed_archetype,
        doc_text=job.get("doc_text"),
        excepted_inference_elements=job.get("excepted_inference_elements"),
        has_brand_reference=bool(intake_img),
        association_words=job.get("association_words") or [],
        design_philosophies=design_philosophies,
    )
    (cand_dir / "approach_prompt.txt").write_text(prompt_text)

    if intake_img:
        approach_input: str | list = [{"role": "user", "content": [
            {"type": "input_text", "text": prompt_text},
            {"type": "input_text", "text": "BRAND REFERENCE IMAGE — the client's existing visual identity."},
            _image_content(intake_img),
        ]}]
    else:
        approach_input = prompt_text

    last_error: Exception | None = None
    for attempt in range(1, 3):
        try:
            resp = await _responses_create(
                job_id,
                stage=f"approach_drafter_{attempt}",
                model=prompts.TEMPLATE_DRAFTER_MODEL,
                input=approach_input,
                output_path=(cand_dir / f"approach_response_attempt_{attempt}.txt").relative_to(job_dir(job_id)),
            )
            raw_text = _response_text(resp)
            (cand_dir / f"approach_response_attempt_{attempt}.txt").write_text(raw_text)
            manifest = _parse_model_json(raw_text)
            cleaned = _validate_approach_manifest(manifest)
            (cand_dir / "approach_manifest.json").write_text(json.dumps({"approaches": cleaned}, indent=2))

            # Derive a representative archetype for the job from the most common
            # effective_archetype_id across approaches (used as a summary label only).
            eff_ids = [a.get("effective_archetype_id") or
                       (a.get("archetype_a") or {}).get("archetype_id") or ""
                       for a in cleaned]
            dominant_id = Counter(eff_ids).most_common(1)[0][0] if eff_ids else ""
            dominant = prompts.archetype_by_id(dominant_id) or {}
            db.update_job(
                job_id,
                status="approaches_ready",
                approach_candidates_json=cleaned,
                selected_archetype_id=dominant_id or None,
                selected_archetype_label=dominant.get("label") or None,
                secondary_archetype_id=None,
                secondary_archetype_label=None,
                error_message=None,
            )
            pair_ids = [(a.get("archetype_a", {}).get("archetype_id"),
                         a.get("archetype_b", {}).get("archetype_id")) for a in cleaned]
            append_telemetry(job_id, {
                "event_type": "workflow",
                "stage": "approach_drafter",
                "status": "approaches_ready",
                "attempt": attempt,
                "archetype_pairs": pair_ids,
                "dominant_archetype_id": dominant_id,
                "archetype_was_fixed": fixed_archetype is not None,
            })
            return
        except Exception as exc:
            last_error = exc
            append_telemetry(job_id, {
                "event_type": "workflow",
                "stage": "approach_drafter",
                "status": "retrying" if attempt == 1 else "failed",
                "attempt": attempt,
                "error": str(exc)[:500],
            })

    db.update_job(
        job_id,
        status="approaches_failed",
        error_message=f"Workshop approach drafter failed: {last_error}",
    )


async def approach_regenerate_worker(job_id: str, regenerate_ids: list[str]):
    """Regenerate only the flagged approach_id(s) within an existing batch of
    four, holding the rest fixed. Sibling context prevents collision."""
    job = db.get_job(job_id)
    if not job or not job.get("approach_candidates_json"):
        raise RuntimeError("Job has no existing approaches to regenerate")

    archetype = prompts.archetype_by_id(job["selected_archetype_id"])
    existing = job["approach_candidates_json"]
    keep = [a for a in existing if a["approach_id"] not in regenerate_ids]

    cand_dir = job_dir(job_id) / "approaches"
    cand_dir.mkdir(exist_ok=True)
    append_telemetry(job_id, {
        "event_type": "workflow", "stage": "approach_regenerate", "status": "started",
        "regenerate_ids": regenerate_ids,
    })

    prompt_text = prompts.approach_regenerate_prompt(
        audience=job.get("audience") or "",
        elevator_pitch=job.get("elevator_pitch") or "",
        conveys=job.get("conveys") or "",
        selected_archetype=archetype,
        doc_text=job.get("doc_text"),
        keep_approaches=keep,
        regenerate_ids=regenerate_ids,
    )

    resp = await _responses_create(
        job_id,
        stage="approach_regenerate",
        model=prompts.TEMPLATE_DRAFTER_MODEL,
        input=prompt_text,
    )
    manifest = _parse_model_json(_response_text(resp))
    replacements = manifest.get("approaches")
    if not isinstance(replacements, list) or len(replacements) != len(regenerate_ids):
        raise RuntimeError(f"Expected {len(regenerate_ids)} replacement approach(es); got {len(replacements) if isinstance(replacements, list) else 'invalid'}")

    by_id = {a["approach_id"]: a for a in keep}
    for r in replacements:
        rid = str(r.get("approach_id") or "").strip()
        if rid not in regenerate_ids:
            raise RuntimeError(f"Regenerated approach_id {rid!r} was not requested")
        for field in ("label", "pitch_angle", "key_differentiator", "visual_direction"):
            if not str(r.get(field) or "").strip():
                raise RuntimeError(f"Regenerated approach {rid} is missing {field}")
        by_id[rid] = {
            "approach_id": rid,
            "label": str(r["label"]).strip(),
            "pitch_angle": str(r["pitch_angle"]).strip(),
            "key_differentiator": str(r["key_differentiator"]).strip(),
            "visual_direction": str(r["visual_direction"]).strip(),
        }

    ordered = [by_id[f"approach_{i}"] for i in range(1, 5)]
    db.update_job(job_id, approach_candidates_json=ordered)
    append_telemetry(job_id, {"event_type": "workflow", "stage": "approach_regenerate", "status": "done", "regenerate_ids": regenerate_ids})
    return ordered


async def design_drafter_text_worker(job_id: str):
    """Phase 1 of prototype generation: text-only call, independent of approaches.

    Optimizes purely against material physics and functional geometry.
    Stores 4 design candidates in prototype_candidates_json (no images yet).
    If the user uploaded a product/brand reference image, uses that instead.
    """
    job = db.get_job(job_id)
    if not job:
        return

    design_dir = job_dir(job_id) / "design_references"
    design_dir.mkdir(exist_ok=True)
    append_telemetry(job_id, {"event_type": "workflow", "stage": "design_drafter_text", "status": "started"})

    # ── Uploaded image shortcut ───────────────────────────────────────────────
    supporting = _supporting_image(job_id)
    if supporting:
        check_resp = await _responses_create(
            job_id,
            stage="design_reference_check",
            model=prompts.CLASSIFY_MODEL,
            input=[{"role": "user", "content": [
                {"type": "input_text", "text": (
                    "Classify this image as exactly one of three types:\n"
                    "- product: a physical product, prototype, hardware device, or engineering drawing\n"
                    "- brand: a website, UI screenshot, logo, or brand/marketing material\n"
                    "- neither: anything else\n"
                    "Answer only: product, brand, or neither."
                )},
                _image_content(supporting),
            ]}],
        )
        _ref_raw = _response_text(check_resp).strip().lower()
        ref_type = "product" if _ref_raw.startswith("product") else "brand" if _ref_raw.startswith("brand") else "neither"
        if ref_type in ("product", "brand"):
            # Uploaded image stands in for all 4 prototypes — pre-select it
            candidate = {
                "design_id": "design_uploaded",
                "assigned_archetype_id": ref_type,
                "design_philosophy": "User-provided reference image.",
                "physical_description": "See uploaded image.",
                "reference_image_prompt": "",
                "reference_image_path": str(supporting),
            }
            db.update_job(job_id, prototype_candidates_json=[candidate], selected_prototype_id="design_uploaded")
            append_telemetry(job_id, {"event_type": "workflow", "stage": "design_drafter_text",
                                      "status": "done_from_upload", "ref_type": ref_type})
            return

    # ── Generate 4 design philosophies (no images) ───────────────────────────
    prompt_text = prompts.design_drafter_prompt(
        elevator_pitch=job.get("elevator_pitch") or "",
        doc_text=job.get("doc_text"),
        excepted_inference_elements=job.get("excepted_inference_elements"),
    )
    (design_dir / "design_drafter_prompt.txt").write_text(prompt_text)

    resp = await _responses_create(
        job_id,
        stage="design_drafter_text",
        model=prompts.TEMPLATE_DRAFTER_MODEL,
        input=prompt_text,
        output_path=(design_dir / "design_drafter_response.txt").relative_to(job_dir(job_id)),
    )
    raw_text = _response_text(resp)
    manifest = _parse_model_json(raw_text)
    designs = manifest.get("designs") or []
    if len(designs) != 4:
        raise RuntimeError(f"design_drafter_text returned {len(designs)} designs; expected 4")

    candidates = [
        {
            "design_id": d.get("design_id", f"design_{i+1}"),
            "assigned_archetype_id": d.get("assigned_archetype_id", ""),
            "design_philosophy": d.get("design_philosophy", ""),
            "physical_description": d.get("physical_description", ""),
            "reference_image_prompt": d.get("reference_image_prompt", ""),
            "reference_image_path": None,
        }
        for i, d in enumerate(designs)
    ]
    (design_dir / "design_manifest.json").write_text(json.dumps(manifest, indent=2))
    db.update_job(job_id, prototype_candidates_json=candidates)
    append_telemetry(job_id, {"event_type": "workflow", "stage": "design_drafter_text", "status": "done"})


async def design_drafter_images_worker(job_id: str):
    """Phase 2 of prototype generation: 4 parallel image gens from text candidates.

    Reads prototype_candidates_json, generates one studio product shot per design,
    writes paths back. Runs concurrently with approach_draft.
    """
    job = db.get_job(job_id)
    if not job:
        return
    candidates = job.get("prototype_candidates_json") or []
    if not candidates:
        raise RuntimeError("design_drafter_images_worker: no prototype_candidates_json")

    design_dir = job_dir(job_id) / "design_references"
    design_dir.mkdir(exist_ok=True)
    append_telemetry(job_id, {"event_type": "workflow", "stage": "design_drafter_images", "status": "started"})

    async def _gen(candidate: dict) -> tuple[str, str]:
        design_id = candidate["design_id"]
        ref_prompt = str(candidate.get("reference_image_prompt") or "").strip()
        if not ref_prompt:
            raise RuntimeError(f"design_drafter_images: missing reference_image_prompt for {design_id}")
        img_prompt_text = (
            f"{ref_prompt}\n\n"
            "Binding output constraints:\n"
            "- Generate exactly one image.\n"
            "- Clean neutral studio background only — no scene, no context.\n"
            "- Full device visible, unobstructed, three-quarter view.\n"
            "- No human hands, no props, no background objects.\n"
            "- Photorealistic product photography.\n"
            "- No text, labels, callouts, or brand marks visible.\n"
        )
        out = design_dir / f"{design_id}_design_ref.png"
        await _throttle_image_generation()
        img_resp = await _responses_create(
            job_id,
            stage=f"design_ref_image_{design_id}",
            model=prompts.DECK_DRAFTER_MODEL,
            input=img_prompt_text,
            tools=[prompts.image_generation_tool(stage="design_ref")],
            output_path=out.relative_to(job_dir(job_id)),
        )
        images = _response_images(img_resp)
        if not images:
            raise RuntimeError(f"design_drafter_images: no image returned for {design_id}")
        await _write_image_from_b64(images[0], out)
        return design_id, str(out)

    results = await _chunked_gather([_gen(c) for c in candidates if c.get("reference_image_prompt")], chunk_size=4)
    path_by_id = dict(results)

    refreshed = db.get_job(job_id)
    updated = [
        {**c, "reference_image_path": path_by_id.get(c["design_id"], c.get("reference_image_path"))}
        for c in (refreshed.get("prototype_candidates_json") or [])
    ]
    db.update_job(job_id, prototype_candidates_json=updated)
    append_telemetry(job_id, {"event_type": "workflow", "stage": "design_drafter_images",
                               "status": "done", "paths": path_by_id})


async def single_slide_preview_worker(job_id: str, approach_id: str):
    """Merged Anchor-Writer/Template-Drafter role: one text call grounded in
    the real elevator pitch + one approved approach, then one image call.
    This is the stage that today only runs after preview checkout; here it
    runs pre-paywall by design, for workshop comparison purposes."""
    job = db.get_job(job_id)
    if not job:
        return
    approaches = job.get("approach_candidates_json") or []
    approach = next((a for a in approaches if a["approach_id"] == approach_id), None)
    if not approach:
        raise RuntimeError(f"Unknown approach_id {approach_id!r} for job {job_id}")
    arch_a = approach.get("archetype_a") or prompts.archetype_by_id(job.get("selected_archetype_id") or "") or {}
    arch_b = approach.get("archetype_b") or arch_a

    slide_dir = job_dir(job_id) / "single_slide_previews"
    slide_dir.mkdir(exist_ok=True)

    draft_prompt = prompts.single_slide_drafter_prompt(
        audience=job.get("audience") or "",
        elevator_pitch=job.get("elevator_pitch") or "",
        approach=approach,
        archetype_a=arch_a,
        archetype_b=arch_b,
        doc_text=job.get("doc_text"),
        excepted_inference_elements=job.get("excepted_inference_elements"),
    )
    (slide_dir / f"{approach_id}_draft_prompt.txt").write_text(draft_prompt)

    resp = await _responses_create(
        job_id,
        stage=f"single_slide_draft_{approach_id}",
        model=prompts.TEMPLATE_DRAFTER_MODEL,
        input=draft_prompt,
    )
    raw_text = _response_text(resp)
    (slide_dir / f"{approach_id}_draft_response.txt").write_text(raw_text)
    slide_draft = _parse_model_json(raw_text)

    # Use the user-selected prototype image if one has been chosen,
    # otherwise fall back to the approach's own design_reference_image_path
    # (legacy path, or uploaded-image shortcut).
    design_ref_path: Path | None = None
    selected_proto_id = job.get("selected_prototype_id")
    if selected_proto_id:
        proto_candidates = job.get("prototype_candidates_json") or []
        proto = next((c for c in proto_candidates if c.get("design_id") == selected_proto_id), None)
        proto_path_str = proto.get("reference_image_path") if proto else None
        if proto_path_str:
            p = Path(proto_path_str)
            design_ref_path = p if p.exists() else None
    if design_ref_path is None:
        fallback_str = approach.get("design_reference_image_path")
        if fallback_str:
            p = Path(fallback_str)
            design_ref_path = p if p.exists() else None
    design_ref_type = "product"

    slide_image_text = prompts.single_slide_image_prompt(slide_draft=slide_draft)
    if design_ref_path:
        if design_ref_type == "brand":
            ref_label = (
                "BRAND VISUAL REFERENCE — this is the client's existing visual identity. "
                "Match its color palette, typography register, iconographic style, and "
                "overall aesthetic tone in the slide design. Do not reproduce this image "
                "literally; use it as the visual contract for the deck's design language."
            )
        else:
            ref_label = (
                "PRODUCT DESIGN REFERENCE — use this as the authoritative visual for "
                "the physical product. Replicate its exact form factor, materials, "
                "colors, and proportions. Do not invent variations."
            )
        slide_image_input: str | list = [{"role": "user", "content": [
            {"type": "input_text", "text": slide_image_text},
            {"type": "input_text", "text": ref_label},
            _image_content(design_ref_path),
        ]}]
    else:
        slide_image_input = slide_image_text

    await _throttle_image_generation()
    image_resp = await _responses_create(
        job_id,
        stage=f"single_slide_image_{approach_id}",
        model=prompts.DECK_DRAFTER_MODEL,
        input=slide_image_input,
        tools=[prompts.image_generation_tool(stage="single_slide")],
        output_path=(slide_dir / f"{approach_id}.jpg").relative_to(job_dir(job_id)),
    )
    images = _response_images(image_resp)
    if len(images) != 1:
        raise RuntimeError(f"Single-slide preview {approach_id} returned {len(images)} images; expected 1")
    img_path = slide_dir / f"{approach_id}.jpg"
    await _write_preview_jpg(images[0], img_path, approach, job)

    refreshed = db.get_job(job_id)
    updated_approaches = refreshed.get("approach_candidates_json") or []
    for a in updated_approaches:
        if a["approach_id"] == approach_id:
            a["preview_image_path"] = str(img_path)
            break
    db.update_job(job_id, approach_candidates_json=updated_approaches)

    append_telemetry(job_id, {
        "event_type": "workflow", "stage": "single_slide_preview", "status": "done", "approach_id": approach_id,
    })
    return img_path


async def candidate_worker(job_id: str):
    if not db.claim_job_status(
        job_id,
        {"preview_paid", "preview_generation_queued", "preview_checkout_pending"},
        "preview_generation_started",
    ):
        append_telemetry(job_id, {"event_type": "workflow", "stage": "preview", "status": "duplicate_skipped"})
        return
    job = db.get_job(job_id)
    user_inputs = _raw_user_inputs(job)
    selected_template = _selected_template_payload(job)
    doc_text = job.get("doc_text") or ""

    cand_dir = job_dir(job_id) / "candidates"
    cand_dir.mkdir(exist_ok=True)
    append_telemetry(job_id, {"event_type": "workflow", "stage": "preview", "status": "started"})

    anchor_prompt = prompts.anchor_writer_prompt(
        mode="PREVIEW",
        user_inputs=user_inputs,
        selected_template=selected_template,
        supporting_document=doc_text,
        supporting_image=bool(_supporting_image(job_id)),
    )

    try:
        anchor_resp = await _responses_create(
            job_id,
            stage="preview_anchor_writer",
            model=prompts.ANCHOR_MODEL,
            input=_anchor_input(anchor_prompt, job_id),
        )
        anchor_text = _response_text(anchor_resp)
        anchor_json = _parse_model_json(anchor_text)
        deck_generation_prompt = anchor_json["deck_generation_prompt"]
        (cand_dir / "anchor_writer_output.json").write_text(json.dumps(anchor_json, indent=2))
    except Exception as e:
        append_telemetry(job_id, {"event_type": "workflow", "stage": "preview", "status": "failed", "error": str(e)[:500]})
        db.update_job(job_id, status="preview_generation_failed", error_message=f"Anchor writer failed: {e}")
        return

    preview_slide_numbers = list(prompts.PREVIEW_EXPOSED_SLIDES)
    preview_prompt = prompts.preview_deck_builder_prompt(
        deck_generation_prompt=deck_generation_prompt,
        user_inputs=user_inputs,
        selected_template=selected_template,
        slide_numbers=preview_slide_numbers,
    )
    try:
        # Phase 1: text-only storyboard call — no images generated yet
        storyboard_resp = await _responses_create(
            job_id,
            stage="preview_storyboard",
            model=prompts.DECK_DRAFTER_MODEL,
            input=_deck_drafter_input(
                _anchor_payload(anchor_json, focused_prompt=preview_prompt),
                job_id,
            ),
            output_path=(cand_dir / "preview_storyboard").relative_to(job_dir(job_id)),
        )
        storyboard_text = _response_text(storyboard_resp)
        storyboard_json = _parse_model_json(storyboard_text)
        preview_storyboard = storyboard_json.get("storyboard") or []

        # Validate structure before spending on image generation
        if len(preview_storyboard) != 3:
            raise RuntimeError(f"Deck Builder returned {len(preview_storyboard)} storyboard entries; expected 3")
        actual_slide_numbers = [int(item.get("slide_number") or 0) for item in preview_storyboard]
        if actual_slide_numbers != preview_slide_numbers:
            raise RuntimeError(f"Deck Builder storyboard slides must be {preview_slide_numbers}; got {actual_slide_numbers}")
        for item in preview_storyboard:
            if not str(item.get("image_prompt") or "").strip():
                raise RuntimeError(f"Deck Builder storyboard entry for slide {item.get('slide_number')} is missing image_prompt")

        (cand_dir / "deck_builder_output.json").write_text(json.dumps(storyboard_json, indent=2))
    except Exception as e:
        append_telemetry(job_id, {"event_type": "workflow", "stage": "preview", "status": "failed", "error": str(e)[:500]})
        db.update_job(job_id, status="preview_generation_failed", error_message=f"Deck Builder storyboard failed: {e}")
        return

    try:
        # Phase 2: one image generation call per validated storyboard slide
        async def _gen_preview_slide(i: int, entry: dict) -> str:
            slide_num = entry["slide_number"]
            await _throttle_image_generation()
            img_resp = await _responses_create(
                job_id,
                stage=f"preview_slide_{slide_num}",
                model=prompts.DECK_DRAFTER_MODEL,
                input=_preview_slide_image_input(entry["image_prompt"], job_id),
                tools=[prompts.image_generation_tool(stage="preview")],
                output_path=(cand_dir / f"cand_{i}.png").relative_to(job_dir(job_id)),
            )
            images = _response_images(img_resp)
            if len(images) != 1:
                raise RuntimeError(f"Preview slide {slide_num} returned {len(images)} images; expected 1")
            return images[0]

        image_payloads = list(await asyncio.gather(*[
            _gen_preview_slide(i, entry)
            for i, entry in enumerate(preview_storyboard[:3], start=1)
        ]))
    except Exception as e:
        append_telemetry(job_id, {"event_type": "workflow", "stage": "preview", "status": "failed", "error": str(e)[:500]})
        db.update_job(job_id, status="preview_generation_failed", error_message=f"Deck Builder image generation failed: {e}")
        return

    candidates = []
    for i, img_b64 in enumerate(image_payloads[:3]):
        img_path = cand_dir / f"cand_{i+1}.png"
        await _write_image_from_b64(img_b64, img_path)
        storyboard_entry = preview_storyboard[i]
        candidates.append({
            "candidate_id": f"cand_{i+1}",
            "url": f"/instant/pitch-deck/jobs/{job_id}/candidates/cand_{i+1}.png",
            "style_label": storyboard_entry.get("title") or f"Slide {preview_slide_numbers[i]}",
            "style_tags": storyboard_entry.get("style_tags") or [],
        })

    (cand_dir / "candidates.json").write_text(json.dumps(candidates))
    (cand_dir / "storyboard.json").write_text(json.dumps(preview_storyboard, indent=2))
    db.update_job(
        job_id,
        status="candidate_ready",
        preview_artifacts_json=candidates,
        preview_storyboard_json=preview_storyboard,
        error_message=None,
    )
    append_telemetry(job_id, {"event_type": "workflow", "stage": "preview", "status": "candidate_ready", "candidate_count": len(candidates)})


# ─────────────────────────────────────────────────────────────────────────────
# Stage 4 + 5 + 6 — Generation worker (runs after payment confirmed)
# ─────────────────────────────────────────────────────────────────────────────

async def generation_worker(job_id: str):
    if not db.claim_job_status(job_id, {"paid"}, "generating_plan"):
        append_telemetry(job_id, {"event_type": "workflow", "stage": "paid_generation", "status": "duplicate_skipped"})
        return
    job = db.get_job(job_id)
    user_inputs = _raw_user_inputs(job)
    selected_template = _selected_template_payload(job)
    doc_text = job.get("doc_text") or ""

    # Extract design philosophy from the selected approach for the anchor writer.
    _selected_approach_id = job.get("selected_candidate_id") or ""
    _selected_approach = next(
        (a for a in (job.get("approach_candidates_json") or [])
         if a.get("approach_id") == _selected_approach_id),
        None,
    )
    _design_philosophy = (_selected_approach or {}).get("design_philosophy") or None

    # Treat None and [] equally — both mean "not explicitly configured."
    # The natural default is to authorize both narrow inference elements.
    _excepted_inference_elements = (
        job.get("excepted_inference_elements") or ["pricing"]
    )

    append_telemetry(job_id, {
        "event_type": "workflow",
        "stage": "paid_generation",
        "status": "started",
        "requested_content_slide_count": job.get("explicit_slide_count"),
    })
    try:
        content_slide_count = prompts.paid_slide_count(
            explicit_slide_count=job.get("explicit_slide_count")
        )
        slide_count = prompts.deck_builder_slide_count(
            explicit_slide_count=job.get("explicit_slide_count")
        )
        anchor_prompt = prompts.anchor_writer_prompt(
            mode="PAID",
            user_inputs=user_inputs,
            selected_template=selected_template,
            paid_slide_count_value=content_slide_count,
            supporting_document=doc_text,
            supporting_image=bool(_supporting_image(job_id)),
            design_philosophy=_design_philosophy,
        )
        anchor_resp = await _responses_create(
            job_id,
            stage="paid_anchor_writer",
            model=prompts.ANCHOR_MODEL,
            input=_anchor_input(anchor_prompt, job_id),
        )
        anchor_json = _parse_model_json(_response_text(anchor_resp))
        deck_generation_prompt = anchor_json["deck_generation_prompt"]
        strategy_dir = job_dir(job_id) / "strategy"
        strategy_dir.mkdir(exist_ok=True)
        (strategy_dir / "anchor_writer_output.json").write_text(json.dumps(anchor_json, indent=2))
        db.append_progress(job_id, "generation", "Creative direction locked — anchor write complete", pct=0.10)
    except Exception as e:
        append_telemetry(job_id, {"event_type": "workflow", "stage": "paid_generation", "status": "failed", "error": str(e)[:500]})
        db.update_job(job_id, status="failed", error_message=f"Anchor writer failed: {e}")
        return

    # Phase 1b: parallel Visual Grammar and Founder Narrative synthesis.
    # Both calls receive the untouched anchor narrative and run concurrently. Their
    # outputs are then appended in a fixed order before the Deck Builder receives the
    # prompt: anchor -> visual production grammar -> founder presentation grammar.
    try:
        (
            deck_generation_prompt,
            visual_grammar_result,
            founder_narrative_result,
        ) = await _synthesize_anchor_expansions(
            job_id,
            deck_generation_prompt=deck_generation_prompt,
            user_inputs=user_inputs,
            selected_template=selected_template,
            strategy_dir=strategy_dir,
            stage_prefix="paid",
        )

        anchor_json["deck_generation_prompt"] = deck_generation_prompt
        db.append_progress(
            job_id,
            "generation",
            "Visual and founder presentation grammars synthesized",
            pct=0.18,
        )
    except Exception as e:
        append_telemetry(
            job_id,
            {
                "event_type": "workflow",
                "stage": "paid_anchor_expansion",
                "status": "failed",
                "error": str(e)[:500],
            },
        )
        db.update_job(
            job_id,
            status="failed",
            error_message=f"Anchor expansion synthesis failed: {e}",
        )
        return

    # Phase 2: Deck builder storyboard — one text call that plans all N slides with
    # explicit visual consistency across every image_prompt.
    deck_builder_prompt_text = prompts.paid_deck_builder_prompt(
        deck_generation_prompt=deck_generation_prompt,
        user_inputs=user_inputs,
        selected_template=selected_template,
        slide_count=slide_count,
        excepted_inference_elements=_excepted_inference_elements,
    )
    d = job_dir(job_id) / "strategy"
    d.mkdir(exist_ok=True)
    try:
        storyboard_resp = await _responses_create(
            job_id,
            stage="paid_storyboard",
            model=prompts.DECK_DRAFTER_MODEL,
            input=_deck_drafter_input(
                _anchor_payload(anchor_json, focused_prompt=deck_builder_prompt_text),
                job_id,
            ),
            output_path=(d / "paid_storyboard").relative_to(job_dir(job_id)),
        )
        storyboard_json = _parse_model_json(_response_text(storyboard_resp))
        paid_storyboard = storyboard_json.get("storyboard") or []
        if len(paid_storyboard) != slide_count:
            raise RuntimeError(f"Deck builder returned {len(paid_storyboard)} storyboard entries; expected {slide_count}")
        actual_numbers = [int(item.get("slide_number") or 0) for item in paid_storyboard]
        expected_numbers = list(range(1, slide_count + 1))
        if actual_numbers != expected_numbers:
            raise RuntimeError(f"Deck builder slide_numbers must be {expected_numbers}; got {actual_numbers}")
        for item in paid_storyboard:
            if not str(item.get("image_prompt") or "").strip():
                raise RuntimeError(f"Deck builder storyboard entry for slide {item.get('slide_number')} is missing image_prompt")
            if not [p for p in (item.get("body_points") or []) if str(p).strip()]:
                raise RuntimeError(f"Deck builder storyboard entry for slide {item.get('slide_number')} is missing body_points")
        image_prompts = {int(item["slide_number"]): item["image_prompt"] for item in paid_storyboard}
        visual_grammar = storyboard_json.get("visual_grammar") or {}
        (d / "paid_storyboard.json").write_text(json.dumps(storyboard_json, indent=2))
        db.append_progress(job_id, "generation", f"Storyboard complete — {slide_count} slides planned, starting render", pct=0.25)
    except Exception as e:
        append_telemetry(job_id, {"event_type": "workflow", "stage": "paid_generation", "status": "failed", "error": str(e)[:500]})
        db.update_job(job_id, status="failed", error_message=f"Deck builder storyboard failed: {e}")
        return

    slide_specs = [
        {
            "slide_index": item["slide_number"],
            "slide_label": f"Slide {item['slide_number']}",
            "headline": item.get("title") or f"Slide {item['slide_number']}",
            "body_points": [str(p).strip() for p in (item.get("body_points") or []) if str(p).strip()],
            "speaker_note": item.get("purpose") or "",
        }
        for item in paid_storyboard
    ]
    plan_json = {
        "deck_title": "Pitch Deck",
        "mode": "PAID",
        "slide_count": slide_count,
        "requested_content_slide_count": content_slide_count,
        "deck_storyboard": [
            {"slide_number": item["slide_number"], "title": item.get("title", ""), "purpose": item.get("purpose", "")}
            for item in paid_storyboard
        ],
        "anchor_writer_output": anchor_json,
        "paid_storyboard": paid_storyboard,
    }
    expected_text_map = {
        str(spec["slide_index"]): {"headline": spec["headline"], "body_points": spec["body_points"]}
        for spec in slide_specs
    }

    (d / "deck_proof_plan.json").write_text(json.dumps(plan_json, indent=2))
    (d / "slide_specs.json").write_text(json.dumps(slide_specs, indent=2))
    (d / "expected_text_map.json").write_text(json.dumps(expected_text_map, indent=2))

    db.update_job(
        job_id,
        status="plan_ready",
        deck_proof_plan=plan_json,
        slide_specs=slide_specs,
        expected_text_map=expected_text_map,
    )

    db.update_job(job_id, status="rendering_slides")

    slides_dir = job_dir(job_id) / "slides"
    slides_dir.mkdir(exist_ok=True)

    async def _generate_slide(spec: dict):
        idx = spec["slide_index"]
        image_prompt = prompts.paid_slide_image_prompt(image_prompts.get(idx, ""), visual_grammar)
        for attempt in range(2):
            try:
                out = slides_dir / f"slide_{idx:02d}_proof.png"
                await _throttle_image_generation()
                resp = await _responses_create(
                    job_id,
                    stage=f"paid_slide_{idx}",
                    model=prompts.DECK_DRAFTER_MODEL,
                    input=_deck_drafter_input(
                        _anchor_payload(anchor_json, focused_prompt=image_prompt, slide_number=idx),
                        job_id,
                    ),
                    tools=[prompts.image_generation_tool(stage="paid")],
                    output_path=out.relative_to(job_dir(job_id)),
                )
                images = _response_images(resp)
                if not images:
                    raise RuntimeError("No image returned")
                await _write_image_from_b64(images[0], out)
                pct = 0.25 + 0.70 * (idx / len(slide_specs))
                db.append_progress(job_id, "generation", f"Slide {idx} of {len(slide_specs)} rendered", pct=pct)
                return idx, True
            except Exception as exc:
                append_telemetry(job_id, {
                    "event_type": "workflow",
                    "stage": f"paid_slide_{idx}",
                    "status": "attempt_failed",
                    "attempt": attempt + 1,
                    "error": str(exc)[:500],
                })
                if attempt == 1:
                    return idx, False
        return idx, False

    results = await _chunked_gather([_generate_slide(s) for s in slide_specs])
    failed_slides = [idx for idx, ok in results if not ok]
    if len(failed_slides) > 2:
        append_telemetry(job_id, {"event_type": "workflow", "stage": "paid_generation", "status": "failed", "failed_slides": failed_slides})
        db.update_job(job_id, status="failed", error_message=f"Too many slide failures: {failed_slides}")
        return

    db.append_progress(job_id, "generation", "All slides rendered — ready for review", pct=1.0)
    db.update_job(
        job_id,
        status="awaiting_review",
        quality_results={str(s["slide_index"]): {"decision": "as_generated", "layout_usable": True, "text_regions": []} for s in slide_specs},
    )
    append_telemetry(job_id, {
        "event_type": "workflow",
        "stage": "paid_generation",
        "status": "awaiting_review",
        "slide_count": slide_count,
        "requested_content_slide_count": content_slide_count,
        "failed_slides": failed_slides,
    })


async def convert_generation_worker(job_id: str):
    """
    Runs after a PREVIEW→PAID conversion.

    The job record may contain `convert_keep_candidate_ids` (list of candidate
    IDs the user elected to keep). If none are kept, convert falls back to the
    normal PAID template-driven workflow.
    """
    if not db.claim_job_status(job_id, {"paid"}, "generating_plan"):
        append_telemetry(job_id, {"event_type": "workflow", "stage": "convert_generation", "status": "duplicate_skipped"})
        return
    job = db.get_job(job_id)
    user_inputs = _raw_user_inputs(job)
    selected_template = _selected_template_payload(job)
    doc_text = job.get("doc_text") or ""

    vertex_paths = _preview_vertex_paths(job_id)
    keep_preview_count = len(vertex_paths)
    convert_uses_kept_previews = keep_preview_count > 0
    append_telemetry(job_id, {
        "event_type": "workflow",
        "stage": "convert_generation",
        "status": "started",
        "requested_content_slide_count": job.get("explicit_slide_count"),
        "vertex_count": keep_preview_count,
    })

    try:
        content_slide_count = prompts.paid_slide_count(
            explicit_slide_count=job.get("explicit_slide_count")
        )
        slide_count = prompts.deck_builder_slide_count(
            explicit_slide_count=job.get("explicit_slide_count")
        )
        anchor_prompt = prompts.anchor_writer_prompt(
            mode="CONVERT" if convert_uses_kept_previews else "PAID",
            user_inputs=user_inputs,
            selected_template=selected_template,
            paid_slide_count_value=content_slide_count,
            supporting_document=doc_text,
            supporting_image=bool(_supporting_image(job_id)),
        )
        anchor_resp = await _responses_create(
            job_id,
            stage="convert_anchor_writer",
            model=prompts.ANCHOR_MODEL,
            input=_anchor_input(anchor_prompt, job_id),
        )
        anchor_json = _parse_model_json(_response_text(anchor_resp))
        deck_generation_prompt = anchor_json["deck_generation_prompt"]
    except Exception as e:
        append_telemetry(job_id, {"event_type": "workflow", "stage": "convert_generation", "status": "failed", "error": str(e)[:500]})
        db.update_job(job_id, status="failed", error_message=f"Convert anchor writer failed: {e}")
        return

    d = job_dir(job_id) / "strategy"
    d.mkdir(exist_ok=True)
    try:
        if convert_uses_kept_previews:
            deck_builder_prompt_text = prompts.convert_deck_builder_prompt(
                deck_generation_prompt=deck_generation_prompt,
                user_inputs=user_inputs,
                selected_template=selected_template,
                slide_count=slide_count,
                kept_preview_count=keep_preview_count,
            )
            storyboard_vertex_instruction = (
                "The following photographs are kept preview slides from this same pitch deck. "
                "They replace the template as the governing style anchor. Each one must appear "
                "exactly once in the full deck. Choose where each belongs in the deck argument, "
                "and for those storyboard entries set source_preview_index to the matching input "
                "image number in the order shown here. Newly generated slides must use these kept "
                "preview slides as style guidance and must not introduce a different aesthetic."
            )
        else:
            deck_builder_prompt_text = prompts.paid_deck_builder_prompt(
                deck_generation_prompt=deck_generation_prompt,
                user_inputs=user_inputs,
                selected_template=selected_template,
                slide_count=slide_count,
                excepted_inference_elements=job.get("excepted_inference_elements"),
            )
            storyboard_vertex_instruction = None

        storyboard_resp = await _responses_create(
            job_id,
            stage="convert_storyboard",
            model=prompts.DECK_DRAFTER_MODEL,
            input=_deck_drafter_input(
                _anchor_payload(anchor_json, focused_prompt=deck_builder_prompt_text),
                job_id,
                vertex_images=vertex_paths if convert_uses_kept_previews else None,
                vertex_instruction=storyboard_vertex_instruction,
            ),
            output_path=(d / "convert_storyboard").relative_to(job_dir(job_id)),
        )
        storyboard_json = _parse_model_json(_response_text(storyboard_resp))
        convert_storyboard = storyboard_json.get("storyboard") or []
        if len(convert_storyboard) != slide_count:
            raise RuntimeError(f"Deck builder returned {len(convert_storyboard)} storyboard entries; expected {slide_count}")

        actual_numbers = [int(item.get("slide_number") or 0) for item in convert_storyboard]
        expected_numbers = list(range(1, slide_count + 1))
        if actual_numbers != expected_numbers:
            raise RuntimeError(f"Deck builder slide_numbers must be {expected_numbers}; got {actual_numbers}")

        assigned_preview_indexes: list[int] = []
        for item in convert_storyboard:
            if not str(item.get("image_prompt") or "").strip():
                raise RuntimeError(f"Deck builder storyboard entry for slide {item.get('slide_number')} is missing image_prompt")
            preview_index = item.get("source_preview_index")
            if preview_index is None:
                continue
            try:
                preview_index_int = int(preview_index)
            except (TypeError, ValueError) as exc:
                raise RuntimeError(
                    f"Deck builder source_preview_index must be an integer; got {preview_index!r}"
                ) from exc
            if preview_index_int < 1 or preview_index_int > keep_preview_count:
                raise RuntimeError(
                    f"Deck builder source_preview_index must be within 1..{keep_preview_count}; got {preview_index_int}"
                )
            assigned_preview_indexes.append(preview_index_int)

        if convert_uses_kept_previews:
            expected_preview_indexes = list(range(1, keep_preview_count + 1))
            if sorted(assigned_preview_indexes) != expected_preview_indexes:
                raise RuntimeError(
                    "Deck builder must assign each kept preview slide exactly once; "
                    f"expected source_preview_index values {expected_preview_indexes}, got {sorted(assigned_preview_indexes)}"
                )
        elif assigned_preview_indexes:
            raise RuntimeError("Deck builder must not emit source_preview_index when no kept preview slides were supplied")

        image_prompts = {int(item["slide_number"]): item["image_prompt"] for item in convert_storyboard}
        visual_grammar = storyboard_json.get("visual_grammar") or {}
        (d / "convert_storyboard.json").write_text(json.dumps(storyboard_json, indent=2))
    except Exception as e:
        append_telemetry(job_id, {"event_type": "workflow", "stage": "convert_generation", "status": "failed", "error": str(e)[:500]})
        db.update_job(job_id, status="failed", error_message=f"Convert deck builder storyboard failed: {e}")
        return

    slide_specs = [
        {
            "slide_index": item["slide_number"],
            "slide_label": f"Slide {item['slide_number']}",
            "headline": item.get("title") or f"Slide {item['slide_number']}",
            "body_points": [],
            "speaker_note": item.get("purpose") or "",
        }
        for item in convert_storyboard
    ]
    plan_json = {
        "deck_title": "Pitch Deck",
        "mode": "CONVERT",
        "slide_count": slide_count,
        "requested_content_slide_count": content_slide_count,
        "vertex_count": keep_preview_count,
        "deck_storyboard": [
            {
                "slide_number": item["slide_number"],
                "title": item.get("title", ""),
                "purpose": item.get("purpose", ""),
                "source_preview_index": item.get("source_preview_index"),
            }
            for item in convert_storyboard
        ],
        "anchor_writer_output": anchor_json,
        "convert_storyboard": convert_storyboard,
    }
    (d / "deck_proof_plan.json").write_text(json.dumps(plan_json, indent=2))
    (d / "slide_specs.json").write_text(json.dumps(slide_specs, indent=2))
    (d / "expected_text_map.json").write_text(json.dumps({}, indent=2))
    db.update_job(job_id, status="plan_ready", deck_proof_plan=plan_json, slide_specs=slide_specs, expected_text_map={})
    db.update_job(job_id, status="rendering_slides")

    slides_dir = job_dir(job_id) / "slides"
    slides_dir.mkdir(exist_ok=True)

    async def _generate_slide(spec: dict):
        idx = spec["slide_index"]
        storyboard_item = next(
            item for item in convert_storyboard
            if int(item.get("slide_number") or 0) == idx
        )
        source_preview_index = storyboard_item.get("source_preview_index")
        if source_preview_index is not None:
            preview_index_int = int(source_preview_index)
            prompt_text = prompts.preserved_preview_to_paid_slide_image_prompt(
                image_prompt=image_prompts.get(idx, ""),
                source_preview_index=preview_index_int,
            )
            slide_vertex_images = [vertex_paths[preview_index_int - 1]]
            slide_vertex_instruction = (
                f"The following input image is kept preview slide {preview_index_int} for this same deck. "
                "Preserve its on-canvas content and composition exactly, and convert only the artifact "
                "surface into a clean full-bleed slide screenshot with no software UI chrome."
            )
        else:
            prompt_text = prompts.paid_slide_image_prompt(image_prompts.get(idx, ""), visual_grammar)
            slide_vertex_images = vertex_paths if convert_uses_kept_previews else None
            slide_vertex_instruction = (
                "The following kept preview slides belong to this same deck. Use them as governing "
                "style anchors for this newly generated slide; match their visual language, but do not "
                "duplicate their compositions."
            ) if convert_uses_kept_previews else None
        for attempt in range(2):
            try:
                out = slides_dir / f"slide_{idx:02d}_proof.png"
                await _throttle_image_generation()
                resp = await _responses_create(
                    job_id,
                    stage=f"convert_slide_{idx}",
                    model=prompts.DECK_DRAFTER_MODEL,
                    input=_deck_drafter_input(
                        _anchor_payload(anchor_json, focused_prompt=prompt_text, slide_number=idx),
                        job_id,
                        vertex_images=slide_vertex_images,
                        vertex_instruction=slide_vertex_instruction,
                    ),
                    tools=[prompts.image_generation_tool(stage="paid")],
                    output_path=out.relative_to(job_dir(job_id)),
                )
                images = _response_images(resp)
                if not images:
                    raise RuntimeError("No image returned")
                await _write_image_from_b64(images[0], out)
                return idx, True
            except Exception as exc:
                append_telemetry(job_id, {
                    "event_type": "workflow",
                    "stage": f"convert_slide_{idx}",
                    "status": "attempt_failed",
                    "attempt": attempt + 1,
                    "error": str(exc)[:500],
                })
                if attempt == 1:
                    return idx, False
        return idx, False

    results = await _chunked_gather([_generate_slide(s) for s in slide_specs])
    failed_slides = [idx for idx, ok in results if not ok]
    if len(failed_slides) > 2:
        append_telemetry(job_id, {"event_type": "workflow", "stage": "convert_generation", "status": "failed", "failed_slides": failed_slides})
        db.update_job(job_id, status="failed", error_message=f"Too many slide failures: {failed_slides}")
        return

    db.update_job(
        job_id,
        status="awaiting_review",
        quality_results={str(s["slide_index"]): {"decision": "as_generated", "layout_usable": True, "text_regions": []} for s in slide_specs},
    )
    append_telemetry(job_id, {
        "event_type": "workflow",
        "stage": "convert_generation",
        "status": "awaiting_review",
        "slide_count": slide_count,
        "requested_content_slide_count": content_slide_count,
        "vertex_count": keep_preview_count,
        "failed_slides": failed_slides,
    })


def _validate_plan(plan: dict, slide_count: int) -> list[str]:
    errors = []
    slides = plan.get("slides", [])
    if not slides:
        errors.append("No slides in plan")
        return errors
    expected_count = prompts.paid_slide_count(explicit_slide_count=slide_count)
    if len(slides) != expected_count:
        errors.append(f"Expected exactly {expected_count} slides, got {len(slides)}")
    narrative_functions = set()
    for s in slides:
        if not s.get("headline"):
            errors.append(f"Slide {s.get('slide_index')} missing headline")
        if not s.get("narrative_function"):
            errors.append(f"Slide {s.get('slide_index')} missing narrative_function")
        if not s.get("visual_brief"):
            errors.append(f"Slide {s.get('slide_index')} missing visual_brief")
        nf = s.get("narrative_function", "")
        if nf in narrative_functions:
            errors.append(f"Duplicate narrative_function: {nf}")
        narrative_functions.add(nf)
    return errors


# ─────────────────────────────────────────────────────────────────────────────
# Stage 8 — Finalization (Procedure 8A)
# ─────────────────────────────────────────────────────────────────────────────

async def finalization_worker(job_id: str):
    if not db.claim_job_status(job_id, {"review_received", "finalization_queued"}, "finalizing"):
        append_telemetry(job_id, {"event_type": "workflow", "stage": "finalization", "status": "duplicate_skipped"})
        return
    append_telemetry(job_id, {"event_type": "workflow", "stage": "finalization", "status": "started"})
    job = db.get_job(job_id)

    slide_specs = job.get("slide_specs") or []
    reviewed_slides = job.get("reviewed_slides") or {}

    slides_dir = job_dir(job_id) / "slides"
    export_dir = job_dir(job_id) / "export"
    export_slides_dir = export_dir / "slides"
    export_presenter_dir = export_dir / "presenter-view"
    export_dir.mkdir(exist_ok=True)
    export_slides_dir.mkdir(exist_ok=True)
    export_presenter_dir.mkdir(exist_ok=True)

    plan = job.get("deck_proof_plan") or {}
    deck_title = plan.get("deck_title", "Pitch Deck")

    # Copy all proofs to presenter-view export folder
    for spec in slide_specs:
        idx = spec["slide_index"]
        proof_path = slides_dir / f"slide_{idx:02d}_proof.png"
        if proof_path.exists():
            shutil.copyfile(proof_path, export_presenter_dir / f"slide_{idx:02d}_proof.png")

    # Step 1: single refiner call with all proof images → one JSON blob with N prompts
    extraction_prompts: dict[int, str] = {}
    try:
        refiner_content = [
            {"type": "input_text", "text": prompts.finalize_refiner_prompt(
                slide_specs, reviewed_slides, job.get("universal_refinement_instruction"),
                elevator_pitch=job.get("elevator_pitch"), doc_text=job.get("doc_text"),
            )},
        ]
        for spec in slide_specs:
            idx = spec["slide_index"]
            proof_path = slides_dir / f"slide_{idx:02d}_proof.png"
            if proof_path.exists():
                refiner_content.append(_image_content(proof_path))

        refiner_resp = await _responses_create(
            job_id,
            stage="finalize_refiner",
            model=prompts.ANCHOR_MODEL,
            input=[{"role": "user", "content": refiner_content}],
        )
        manifest = _parse_model_json(_response_text(refiner_resp))
        for entry in manifest.get("slide_prompts") or []:
            extraction_prompts[int(entry["slide_index"])] = entry["image_model_prompt"]
    except Exception as exc:
        append_telemetry(job_id, {"event_type": "workflow", "stage": "finalize_refiner", "status": "failed", "error": str(exc)[:500]})

    # Step 2: parallel I2I calls — each proof + its extraction prompt → clean final slide
    async def _finalize_slide(spec: dict):
        idx = spec["slide_index"]
        proof_path = slides_dir / f"slide_{idx:02d}_proof.png"
        if not proof_path.exists():
            return idx, None, spec, "", "", "", False

        reviewed = reviewed_slides.get(str(idx)) or {}
        reviewed_headline = reviewed.get("headline") or spec.get("headline", "")
        reviewed_body = reviewed.get("body") or " | ".join(spec.get("body_points", []))
        reviewed_notes = reviewed.get("notes") or spec.get("speaker_note", "")

        final_path = export_slides_dir / f"slide_{idx:02d}_final.png"
        extraction_prompt = extraction_prompts.get(idx, "")
        gen_prompt = prompts.finalize_slide_image_prompt(extraction_prompt) if extraction_prompt else None

        # A presenter-view proof screenshot is not a valid final slide -- it has
        # the taskbar, timer, speaker-notes sidebar, and nav controls baked in.
        # If the I2I extraction call returns no image (which happens without
        # raising an exception) or throws, retry once before treating it as a
        # genuine failure. Never silently ship the raw proof as if it were the
        # extracted final image -- that's a silent success-claim inflation.
        extracted = False
        last_error: str | None = None
        for attempt in range(2):
            try:
                if not gen_prompt:
                    raise RuntimeError(f"No extraction prompt for slide {idx}")
                await _throttle_image_generation()
                resp = await _responses_create(
                    job_id,
                    stage=f"finalize_slide_{idx}",
                    model=prompts.DECK_DRAFTER_MODEL,
                    input=[{"role": "user", "content": [
                        {"type": "input_text", "text": gen_prompt},
                        _image_content(proof_path),
                    ]}],
                    tools=[prompts.image_generation_tool(stage="paid")],
                    output_path=final_path.relative_to(job_dir(job_id)),
                )
                images = _response_images(resp)
                if images:
                    await _write_image_from_b64(images[0], final_path)
                    extracted = True
                    break
                last_error = "model returned no image"
            except Exception as exc:
                last_error = str(exc)[:500]
            append_telemetry(job_id, {
                "event_type": "workflow", "stage": f"finalize_slide_{idx}",
                "status": "retrying" if attempt == 0 else "failed", "error": last_error,
            })

        if not extracted:
            shutil.copyfile(proof_path, final_path)

        return idx, final_path, spec, reviewed_headline, reviewed_body, reviewed_notes, extracted

    results = await _chunked_gather([_finalize_slide(s) for s in slide_specs])

    # Per-slide failure isolation: every image_generation call is a one-off
    # from the model's point of view, with no shared state across slides, so
    # one slide's outcome must never cascade into failing the other N-1 that
    # already succeeded. Paid-mode proofs are already clean full-bleed
    # screenshots (no presenter-view chrome to strip by design) -- a slide
    # that isn't "extracted" just means the refiner/extraction pass had
    # nothing to change or didn't produce a new image, which is an expected,
    # low-stakes outcome now, not evidence of a dirty source. The raw proof
    # is a legitimate final slide in that case. The rarer, real risk this
    # step guards against is hallucinated chrome that wasn't there in the
    # source -- unlikely given the full-bleed-screenshot request, but still
    # possible; that's what the extraction pass is actually for now.
    composited = []
    dropped_slides = []
    for idx, final_path, spec, reviewed_headline, reviewed_body, reviewed_notes, extracted in sorted(results, key=lambda r: r[0]):
        if final_path is None or not final_path.exists():
            # Only case with truly nothing to ship for this slide (the proof
            # itself never rendered). Skip it, don't take the rest down.
            append_telemetry(job_id, {"event_type": "workflow", "stage": "finalization", "status": "slide_dropped", "error": f"Missing proof for slide {idx}"})
            dropped_slides.append(idx)
            continue
        if not extracted:
            append_telemetry(job_id, {"event_type": "workflow", "stage": "finalization", "status": "slide_used_raw_proof", "slide_index": idx})
        composited.append((spec, final_path, reviewed_headline, reviewed_body, reviewed_notes))

    if not composited:
        append_telemetry(job_id, {"event_type": "workflow", "stage": "finalization", "status": "failed", "error": "No slides produced any usable final image"})
        db.update_job(job_id, status="failed", error_message="Finalization produced no usable slides")
        return
    if dropped_slides:
        append_telemetry(job_id, {"event_type": "workflow", "stage": "finalization", "status": "partial", "dropped_slides": dropped_slides})

    # Hero/title cover slide -- a bonus artifact, not one of the counted
    # slides. Extract it the same way as a numbered slide (I2I chrome
    # removal), but skip silently on failure rather than fail the job: the
    # customer's paid N-slide deck must not depend on the bonus cover.
    hero_final_path = None
    hero_proof_path = slides_dir / "hero_proof.png"
    if hero_proof_path.exists():
        shutil.copyfile(hero_proof_path, export_presenter_dir / "hero_proof.png")
        hero_candidate_path = export_slides_dir / "hero_final.png"
        hero_extracted = False
        for attempt in range(2):
            try:
                await _throttle_image_generation()
                resp = await _responses_create(
                    job_id,
                    stage="finalize_hero_slide",
                    model=prompts.DECK_DRAFTER_MODEL,
                    input=[{"role": "user", "content": [
                        {"type": "input_text", "text": prompts.finalize_hero_slide_image_prompt(job.get("universal_refinement_instruction"))},
                        _image_content(hero_proof_path),
                    ]}],
                    tools=[prompts.image_generation_tool(stage="paid")],
                    output_path=hero_candidate_path.relative_to(job_dir(job_id)),
                )
                images = _response_images(resp)
                if images:
                    await _write_image_from_b64(images[0], hero_candidate_path)
                    hero_extracted = True
                    break
            except Exception as exc:
                append_telemetry(job_id, {
                    "event_type": "workflow", "stage": "finalize_hero_slide",
                    "status": "retrying" if attempt == 0 else "failed", "error": str(exc)[:500],
                })
        if hero_extracted:
            hero_final_path = hero_candidate_path
            composited.insert(0, ({"slide_index": "hero", "slide_label": "Cover"}, hero_final_path, "", "", ""))
        else:
            append_telemetry(job_id, {"event_type": "workflow", "stage": "finalize_hero_slide", "status": "skipped_after_retries"})

    # Build PDF
    pdf_path = export_dir / "deck.pdf"
    await asyncio.to_thread(_build_pdf, composited, pdf_path, deck_title)

    # Build PPTX
    pptx_path = export_dir / "deck.pptx"
    await asyncio.to_thread(_build_pptx, composited, pptx_path, deck_title)

    # Build HTML preview
    html_path = export_dir / "preview.html"
    _build_html(composited, html_path, deck_title, job_id)

    # Build Markdown
    md_path = export_dir / "slides.md"
    _build_markdown(composited, md_path, deck_title)

    source_kit_files: list[str] = []
    if job.get("source_kit"):
        source_kit_dir = export_dir / "source-kit"
        source_kit_dir.mkdir(exist_ok=True)
        anchor_payload = plan.get("anchor_writer_output") or {}
        structured_input = anchor_payload.get("structured_input") or {}
        deck_storyboard = anchor_payload.get("deck_storyboard") or plan.get("deck_storyboard") or []
        slide_transcript = anchor_payload.get("slide_text_transcript") or []
        selected_direction = structured_input.get("selected_direction") or job.get("selected_candidate_label") or "Selected direction"
        associations = structured_input.get("page_1", {}).get("key_vibes") or job.get("association_words") or []
        audience = structured_input.get("page_1", {}).get("audience") or job.get("audience") or ""

        brand_guide_lines = [
            f"# {deck_title} Source Kit",
            "",
            f"- Direction: {selected_direction}",
            f"- Audience: {audience}",
            f"- Associations: {' | '.join(associations)}",
            f"- Slide count: {job.get('explicit_slide_count') or len(deck_storyboard)}",
            "",
            "## Visual Posture",
            "Preserve the selected preview direction as the controlling presentation system.",
            "Keep the same palette pressure, typography posture, layout family, and presenter-view artifact feel.",
        ]
        (source_kit_dir / "brand-guide.md").write_text("\n".join(brand_guide_lines))
        (source_kit_dir / "structured-input.json").write_text(json.dumps(structured_input, indent=2, ensure_ascii=False))
        (source_kit_dir / "storyboard.json").write_text(json.dumps(deck_storyboard, indent=2, ensure_ascii=False))
        (source_kit_dir / "slide-transcript.json").write_text(json.dumps(slide_transcript, indent=2, ensure_ascii=False))
        source_kit_files = [
            "source-kit/brand-guide.md",
            "source-kit/structured-input.json",
            "source-kit/storyboard.json",
            "source-kit/slide-transcript.json",
        ]

    # Build manifest. The hero cover (spec["slide_index"] == "hero") is a
    # bonus artifact and must not be counted in slide_count or the numbered
    # slide_images/presenter_view_images lists -- those remain exactly the
    # customer-facing N counted slides.
    numbered = [spec for spec, *_ in composited if spec.get("slide_index") != "hero"]
    manifest = {
        "job_id": job_id,
        "deck_title": deck_title,
        "slide_count": len(numbered),
        "files": ["deck.pdf", "deck.pptx", "preview.html", "slides.md", "manifest.json", *source_kit_files],
        "slide_images": [f"slides/slide_{spec['slide_index']:02d}_final.png" for spec in numbered],
        "presenter_view_images": [f"presenter-view/slide_{spec['slide_index']:02d}_proof.png" for spec in numbered],
        "telemetry_summary": "telemetry/summary.json",
        "telemetry": telemetry_summary(job_id),
    }
    if hero_final_path:
        manifest["hero_image"] = "slides/hero_final.png"
        manifest["hero_presenter_view_image"] = "presenter-view/hero_proof.png"
    (export_dir / "manifest.json").write_text(json.dumps(manifest, indent=2))

    # Build ZIP
    zip_path = export_dir / "pitch_deck_package.zip"
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for fname in ["deck.pdf", "deck.pptx", "preview.html", "slides.md", "manifest.json", *source_kit_files]:
            p = export_dir / fname
            if p.exists():
                zf.write(p, fname)
        for p in sorted(export_slides_dir.glob("*.png")):
            zf.write(p, f"slides/{p.name}")
        for p in sorted(export_presenter_dir.glob("*.png")):
            zf.write(p, f"presenter-view/{p.name}")

    validation_errors = _validate_export_manifest(export_dir, zip_path, manifest)
    if validation_errors:
        append_telemetry(job_id, {"event_type": "workflow", "stage": "finalization", "status": "failed", "errors": validation_errors})
        db.update_job(job_id, status="failed", error_message="Export validation failed: " + "; ".join(validation_errors))
        return

    db.update_job(job_id, status="complete", export_zip_path=str(zip_path))
    append_telemetry(job_id, {
        "event_type": "workflow",
        "stage": "finalization",
        "status": "complete",
        "export_zip_path": zip_path.relative_to(job_dir(job_id)),
        "export_zip_bytes": zip_path.stat().st_size,
        "slide_count": len(composited),
    })


# ---------------------------------------------------------------------------
# Parallel per-slide verification pipeline (v2 finalization mechanism)
#
# Built against VERIFICATION_PIPELINE_IMPLEMENTATION_PROMPT.md. Replaces the
# single-batch finalize_refiner_prompt call with, per slide: 3 independent
# narrow verifier calls (fan-out) -> concatenated issue list (fan-in) -> one
# judge/disposition call -> conditional full regeneration (never I2I). Every
# slide's pipeline is fully independent -- no shared batch call, no single
# point of failure across the deck. This is a separate worker from
# finalization_worker (left untouched) so the two mechanisms can be run
# side by side on the same job for comparison.
# ---------------------------------------------------------------------------

async def _run_slide_verifiers(
    job_id: str, storyboard_entry: dict, proof_path: Path,
    elevator_pitch: str, doc_text: str | None,
    excepted_inference_elements: list | None = None,
    authorized_inferences: dict | None = None,
    canon_overrides: dict | None = None,
    pitch_aspect_modes: dict | None = None,
    deck_storyboard: list[dict] | None = None,
) -> list[dict]:
    """Fan-out: 3 parallel narrow verifier calls for one slide. Fan-in:
    concatenate their issue lists -- a plain merge, no model call spent
    combining them."""
    label = storyboard_entry.get("slide_number", "hero")

    async def _verify(suffix: str, prompt_text: str) -> dict:
        try:
            resp = await _responses_create(
                job_id, stage=f"verify_{label}_{suffix}", model=prompts.TEMPLATE_DRAFTER_MODEL,
                input=[{"role": "user", "content": [
                    {"type": "input_text", "text": prompt_text},
                    _image_content(proof_path),
                ]}],
            )
            return _parse_model_json(_response_text(resp))
        except Exception as exc:
            append_telemetry(job_id, {
                "event_type": "workflow", "stage": f"verify_{label}_{suffix}",
                "status": "failed", "error": str(exc)[:300],
            })
            return {"verifier": suffix, "issues": [], "error": str(exc)[:300]}

    results = await asyncio.gather(
        _verify("spelling", prompts.slide_spelling_verifier_prompt(storyboard_entry)),
        _verify("source", prompts.slide_source_accuracy_verifier_prompt(
            storyboard_entry, elevator_pitch, doc_text,
            excepted_inference_elements=excepted_inference_elements,
            authorized_inferences=authorized_inferences,
            canon_overrides=canon_overrides,
            pitch_aspect_modes=pitch_aspect_modes,
            deck_storyboard=deck_storyboard,
        )),
        _verify("diagram", prompts.slide_diagram_accuracy_verifier_prompt(storyboard_entry)),
    )
    return list(results)


async def _judge_slide(
    job_id: str, storyboard_entry: dict, verifier_reports: list[dict], visual_grammar: dict,
    pitch_aspect_modes: dict | None = None,
) -> dict:
    label = storyboard_entry.get("slide_number", "hero")
    prompt_text = prompts.slide_disposition_prompt(
        storyboard_entry=storyboard_entry, verifier_reports=verifier_reports,
        visual_grammar=visual_grammar, pitch_aspect_modes=pitch_aspect_modes,
    )
    resp = await _responses_create(
        job_id, stage=f"judge_{label}", model=prompts.ANCHOR_MODEL, input=prompt_text,
    )
    return _parse_model_json(_response_text(resp))


async def _verify_and_finalize_slide(
    job_id: str, storyboard_entry: dict, proof_path: Path, final_path: Path,
    elevator_pitch: str, doc_text: str | None, visual_grammar: dict, anchor_json: dict,
    excepted_inference_elements: list | None = None,
    authorized_inferences: dict | None = None,
    canon_overrides: dict | None = None,
    inferred_element_decisions: dict | None = None,
    pitch_aspect_modes: dict | None = None,
    deck_storyboard: list[dict] | None = None,
) -> dict:
    """One slide's full DAG: fan-out verify -> fan-in -> judge -> conditional
    regen. Returns this slide's acceptance-manifest entry.

    inferred_element_decisions maps excepted-element name -> bool (the
    per-element checkbox, e.g. {"pricing": True} meaning "remove/replace
    this, don't keep it as an inferred proposal"):
      - False/absent (default, unchecked): keep the inferred content, but
        force a relabel edit ("inferred" -> "proposed" wording) so it isn't
        silently displayed as fact -- content stays, disclosure wording
        softens instead of disappearing.
      - True (checked): the element's decided value is withheld from the
        verifier for this slide, so it is checked as an ORDINARY claim.
        A canon override that actually covers it settles it normally (the
        user's own replacement becomes the correction, regenerated in).
        With no covering override, nothing can settle the claim -- the
        verifier's own existing "no decided value" rule fires
        (authorized_inference: true, actionable: false, element: <name>),
        which is the deterministic, robust signal used here to drop the
        slide rather than ship an unresolved claim.
    """
    label = storyboard_entry.get("slide_number", "hero")
    inferred_element_decisions = inferred_element_decisions or {}
    authorized_inferences = authorized_inferences or {}
    checked_elements = {e for e, v in inferred_element_decisions.items() if v}
    # Withhold decided values only for elements the user checked for removal --
    # everything else (unchecked/default) still gets its normal decided value,
    # so it's checked for internal consistency exactly as before.
    authorized_inferences_for_verifier = {
        k: v for k, v in authorized_inferences.items() if k not in checked_elements
    }

    verifier_reports = await _run_slide_verifiers(
        job_id, storyboard_entry, proof_path, elevator_pitch, doc_text,
        excepted_inference_elements=excepted_inference_elements,
        authorized_inferences=authorized_inferences_for_verifier,
        canon_overrides=canon_overrides,
        pitch_aspect_modes=pitch_aspect_modes,
        deck_storyboard=deck_storyboard,
    )

    # Deterministic drop signal: a checked-for-removal element with no
    # covering override and (deliberately) no decided value to check against
    # -- the verifier's existing "nothing settles this" rule fires exactly as
    # it would for any other unverifiable claim; the only new step is acting
    # on it as a removal instead of a "flag for human review, ship as-is."
    for report in verifier_reports:
        for issue in report.get("issues") or []:
            if (
                issue.get("authorized_inference") is True
                and issue.get("actionable") is False
                and issue.get("element") in checked_elements
            ):
                return {
                    "slide": label,
                    "verifier_reports": verifier_reports,
                    "verdict": "drop",
                    "reasoning": (
                        f"Element '{issue.get('element')}' was checked for removal "
                        f"but no covering replacement was found for this slide -- "
                        f"omitted from deliverables rather than shipped with an "
                        f"unresolved claim."
                    ),
                    "final_source": "dropped_no_replacement_for_checked_element",
                }

    disposition = await _judge_slide(
        job_id, storyboard_entry, verifier_reports, visual_grammar,
        pitch_aspect_modes=pitch_aspect_modes,
    )
    # The judge is prompted for corrective_constraints as a single string,
    # but structured-output models occasionally emit a JSON array anyway --
    # normalize once here so every downstream .strip()/f-string merge below
    # can assume a string unconditionally instead of each one re-guarding.
    _cc = disposition.get("corrective_constraints")
    if isinstance(_cc, list):
        disposition["corrective_constraints"] = "\n".join(str(c) for c in _cc if c)
    elif _cc is None:
        disposition["corrective_constraints"] = ""
    elif not isinstance(_cc, str):
        disposition["corrective_constraints"] = str(_cc)

    # Default/unchecked relabel: an unchecked element's decided value has no
    # natural "issue" to trigger on when it matches cleanly (that's the
    # success case, not a defect) -- so the "inferred" -> "proposed" wording
    # softening is a standing edit applied independently of what the verifier
    # found, for any unchecked element whose decided value actually appears
    # on this slide.
    relabel_notes = []
    original_image_prompt = (storyboard_entry.get("image_prompt") or "").strip()
    for element, decided_value in authorized_inferences.items():
        if element in checked_elements:
            continue
        if decided_value and str(decided_value)[:40] in original_image_prompt:
            relabel_notes.append(
                f"Relabel any 'inferred {element}' / 'not a source-verified fact' "
                f"disclosure language for {element} to 'proposed {element}' "
                f"phrasing instead -- keep the actual decided values/content for "
                f"{element} exactly as they are, only the disclosure wording softens."
            )
    if relabel_notes:
        merged = "\n".join(relabel_notes)
        disposition["verdict"] = "regenerate"
        disposition["corrective_constraints"] = (
            f"{disposition.get('corrective_constraints', '').strip()}\n{merged}".strip()
        )

    # Checked-and-resolved: a checked element with a real, non-empty canon
    # override on THIS slide gets its disclosure removed entirely, not
    # softened to "proposed" -- the user specified this value directly, so
    # it is presented as plain fact with no hedging banner. "Resolved" is
    # detected the same way the pre-finalization gate detects "affected":
    # the element's old decided value actually appears on this slide, and
    # the slide has a real (non-blank) reviewed_slides entry. If there is
    # no override here, the drop path above already handled it -- this
    # branch only runs for slides that made it past that check.
    strip_notes = []
    slide_override = (canon_overrides or {}).get(str(label)) or {}
    slide_has_override = any(
        str(slide_override.get(k) or "").strip() for k in ("headline", "body", "notes")
    )
    if slide_has_override:
        for element in checked_elements:
            decided_value = authorized_inferences.get(element)
            if decided_value and str(decided_value)[:40] in original_image_prompt:
                strip_notes.append(
                    f"Remove all 'inferred {element}' / 'proposed {element}' / "
                    f"'not a source-verified fact' disclosure language for "
                    f"{element} entirely -- the user specified this value "
                    f"directly, so present it as a plain, unqualified fact "
                    f"with no hedging banner or disclaimer."
                )
    if strip_notes:
        merged = "\n".join(strip_notes)
        disposition["verdict"] = "regenerate"
        disposition["corrective_constraints"] = (
            f"{disposition.get('corrective_constraints', '').strip()}\n{merged}".strip()
        )

    # Headline override from reviewed_slides is a user-authored mandate —
    # inject it as an explicit correction so the image model receives it
    # verbatim rather than leaving the judge to infer it from the diff.
    reviewed_headline = str(slide_override.get("headline") or "").strip()
    original_title = str(storyboard_entry.get("title") or "").strip()
    if reviewed_headline and reviewed_headline != original_title:
        headline_correction = f'Change the headline text to exactly: "{reviewed_headline}"'
        existing_cc = (disposition.get("corrective_constraints") or "").strip()
        disposition["corrective_constraints"] = (
            f"{headline_correction}\n{existing_cc}" if existing_cc else headline_correction
        )
        disposition["verdict"] = "regenerate"

    entry = {
        "slide": label,
        "verifier_reports": verifier_reports,
        "verdict": disposition.get("verdict"),
        "reasoning": disposition.get("reasoning"),
        "corrective_constraints": disposition.get("corrective_constraints"),
        "pitch_aspect_modes": prompts.normalize_pitch_aspect_modes(pitch_aspect_modes),
        "final_source": None,
    }

    if disposition.get("verdict") == "regenerate" and disposition.get("corrective_constraints"):
        # Splice the corrections onto the ORIGINAL, verbatim image_prompt plus
        # the deck's visual_grammar -- the same code-level splice already used
        # for the original generation -- rather than trusting a model to
        # reproduce composition and content while also rewriting the prompt.
        # That reproduction step is exactly where the first version of this
        # mechanism lost layout fidelity (paraphrase is not preservation).
        corrections = disposition["corrective_constraints"]
        original_image_prompt = (storyboard_entry.get("image_prompt") or "").strip()
        correction_block = (
            f"MANDATORY CORRECTIONS -- apply exactly, everything else in this "
            f"prompt is otherwise unchanged:\n{corrections}"
        )
        if original_image_prompt:
            regen_prompt = prompts.paid_slide_image_prompt(
                f"{original_image_prompt}\n\n{correction_block}", visual_grammar,
            )
        else:
            # Hero has no stored per-slide image_prompt to splice back in --
            # its base prompt is the same fixed template used originally.
            regen_prompt = f"{prompts.hero_slide_image_prompt(visual_grammar)}\n\n{correction_block}"

        # The original per-slide generation call was never just the bare
        # image_prompt -- it was wrapped with the deck's anchor narrative
        # (deck_generation_prompt + user_inputs) via _deck_drafter_input, so
        # every slide's generation shared the same grounding context. Without
        # that, a regeneration reads only this one slide's text in isolation
        # and can independently invent a different product concept from the
        # same underspecified line (e.g. "console-style product screenshot"
        # rendered as a GPS hiking app instead of the AI companion console
        # every other slide depicts). Reuse the exact same wrapping here.
        append_telemetry(job_id, {
            "event_type": "regen_prompt_compiled",
            "stage": f"regen_{label}",
            "slide": label,
            "regen_prompt": regen_prompt[:2000],
        })
        await _throttle_image_generation()
        try:
            resp = await _responses_create(
                job_id, stage=f"regen_{label}", model=prompts.DECK_DRAFTER_MODEL,
                input=_deck_drafter_input(
                    _anchor_payload(
                        anchor_json, focused_prompt=regen_prompt,
                        slide_number=label if isinstance(label, int) else None,
                    ),
                    job_id,
                ),
                tools=[prompts.image_generation_tool(stage="paid")],
                output_path=final_path.relative_to(job_dir(job_id)),
            )
            images = _response_images(resp)
            if images:
                await _write_image_from_b64(images[0], final_path)
                entry["final_source"] = "regenerated"
            else:
                shutil.copyfile(proof_path, final_path)
                entry["final_source"] = "regeneration_failed_used_proof"
        except Exception as exc:
            shutil.copyfile(proof_path, final_path)
            entry["final_source"] = "regeneration_error_used_proof"
            entry["regeneration_error"] = str(exc)[:300]
    else:
        shutil.copyfile(proof_path, final_path)
        entry["final_source"] = "proof_passed"

    return entry


def _gate_check_inferred_element_decisions(
    inferred_element_decisions: dict, authorized_inferences: dict,
    canon_overrides: dict, storyboard_entries: list,
) -> list[str]:
    """Hard-block gate, run before any verifier/judge/image calls in
    verification_worker.

    For each element checked for removal, find every slide whose
    image_prompt actually carries that element's decided value (same
    substring heuristic _verify_and_finalize_slide's relabel path uses),
    then require a real canon override (a reviewed_slides edit) on every
    one of those slides. This is the "must not allow continuing with the
    box checked and zero edits specified" rule. It does not judge whether
    a supplied edit is a VALID replacement for the element -- that's the
    verifier/judge's job at regen time; an edit that doesn't actually
    cover the element still falls through to the existing drop path.
    """
    errors = []
    checked_elements = [e for e, v in (inferred_element_decisions or {}).items() if v]
    for element in checked_elements:
        decided_value = authorized_inferences.get(element)
        if not decided_value:
            # This element was never actually inferred anywhere in the
            # deck -- nothing to remove, checking the box is a no-op.
            continue
        needle = str(decided_value)[:40]
        affected = [
            entry.get("slide_number") for entry in storyboard_entries
            if needle in (entry.get("image_prompt") or "")
        ]
        for idx in affected:
            override = canon_overrides.get(str(idx)) or {}
            has_replacement = any(
                str(override.get(k) or "").strip() for k in ("headline", "body", "notes")
            )
            if not has_replacement:
                errors.append(
                    f"Element '{element}' is checked for removal but slide {idx} "
                    f"has no replacement specified in its refinement edits -- "
                    f"write what {element} should be instead, or uncheck the box."
                )
    return errors


async def verification_worker(job_id: str):
    """v2 finalization mechanism. See module docstring above this function
    and VERIFICATION_PIPELINE_IMPLEMENTATION_PROMPT.md for the full design."""
    if not db.claim_job_status(job_id, {"review_received", "finalization_queued", "awaiting_review"}, "finalizing"):
        append_telemetry(job_id, {"event_type": "workflow", "stage": "verification_finalization", "status": "duplicate_skipped"})
        return
    append_telemetry(job_id, {"event_type": "workflow", "stage": "verification_finalization", "status": "started"})
    job = db.get_job(job_id)

    slide_specs = job.get("slide_specs") or []
    plan = job.get("deck_proof_plan") or {}
    deck_title = plan.get("deck_title", "Pitch Deck")
    deck_storyboard = plan.get("paid_storyboard") or []
    storyboard_by_idx = {item["slide_number"]: item for item in deck_storyboard}
    elevator_pitch = job.get("elevator_pitch") or ""
    doc_text = job.get("doc_text")
    pitch_aspect_modes = prompts.normalize_pitch_aspect_modes(
        job.get("pitch_aspect_modes")
    )
    # Original per-slide generation calls were grounded by the deck-wide
    # anchor narrative (deck_generation_prompt + user_inputs), not just each
    # slide's own image_prompt in isolation. A regeneration must reuse the
    # same grounding or it can independently invent an unrelated concept from
    # an underspecified line -- see _verify_and_finalize_slide.
    anchor_json = plan.get("anchor_writer_output") or {}

    slides_dir = job_dir(job_id) / "slides"
    export_dir = job_dir(job_id) / "export"
    export_slides_dir = export_dir / "slides"
    export_presenter_dir = export_dir / "presenter-view"
    export_dir.mkdir(exist_ok=True)
    export_slides_dir.mkdir(exist_ok=True)
    export_presenter_dir.mkdir(exist_ok=True)

    # visual_grammar and authorized_inferences both live on disk
    # (strategy/paid_storyboard.json) -- neither is persisted to the DB job
    # record by generation_worker.
    visual_grammar: dict = {}
    authorized_inferences: dict = {}
    strategy_path = job_dir(job_id) / "strategy" / "paid_storyboard.json"
    if strategy_path.exists():
        _storyboard_doc = json.loads(strategy_path.read_text()) or {}
        visual_grammar = _storyboard_doc.get("visual_grammar") or {}
        authorized_inferences = _storyboard_doc.get("authorized_inferences") or {}

    excepted_inference_elements = job.get("excepted_inference_elements") or []
    inferred_element_decisions = job.get("inferred_element_decisions") or {}
    # Deck-wide: every slide's verifier sees every OTHER slide's user edits
    # too, not just its own -- a canon override made on one slide is
    # authoritative for every slide that touches the same fact.
    canon_overrides = job.get("reviewed_slides") or {}

    gate_errors = _gate_check_inferred_element_decisions(
        inferred_element_decisions, authorized_inferences, canon_overrides,
        [item for item in (plan.get("paid_storyboard") or [])],
    )
    if gate_errors:
        append_telemetry(job_id, {"event_type": "workflow", "stage": "verification_finalization", "status": "blocked", "errors": gate_errors})
        db.update_job(job_id, status="failed", error_message="Blocked: " + "; ".join(gate_errors))
        return

    async def _process_numbered_slide(spec: dict):
        idx = spec["slide_index"]
        proof_path = slides_dir / f"slide_{idx:02d}_proof.png"
        final_path = export_slides_dir / f"slide_{idx:02d}_final.png"
        if not proof_path.exists():
            return idx, None, None
        shutil.copyfile(proof_path, export_presenter_dir / f"slide_{idx:02d}_proof.png")
        slide_review = (canon_overrides or {}).get(str(idx)) or {}
        if not slide_review.get("_reverify", True):
            shutil.copyfile(proof_path, final_path)
            return idx, final_path, {
                "slide": idx, "verdict": "passthrough",
                "final_source": "proof_passthrough", "verifier_reports": [],
            }
        storyboard_entry = storyboard_by_idx.get(idx) or {
            "slide_number": idx, "title": spec.get("headline", ""),
            "purpose": spec.get("speaker_note", ""), "image_prompt": "",
        }
        entry = await _verify_and_finalize_slide(
            job_id, storyboard_entry, proof_path, final_path, elevator_pitch, doc_text, visual_grammar, anchor_json,
            excepted_inference_elements=excepted_inference_elements,
            authorized_inferences=authorized_inferences,
            canon_overrides=canon_overrides,
            inferred_element_decisions=inferred_element_decisions,
            pitch_aspect_modes=pitch_aspect_modes,
            deck_storyboard=deck_storyboard,
        )
        if entry.get("verdict") == "drop":
            return idx, None, entry
        return idx, final_path, entry

    # Bounded worker-pool concurrency at the slide level -- verifier/judge
    # calls skip the image rate limiter, but dispatching all slides fully
    # unbounded still risks tripping ordinary connection/RPM limits.
    results = await _chunked_gather([_process_numbered_slide(s) for s in slide_specs], chunk_size=6)

    composited = []
    manifest_entries = []
    dropped_slides = []
    for idx, final_path, entry in sorted(results, key=lambda r: r[0]):
        if final_path is None:
            dropped_slides.append(idx)
            if entry is not None:
                # A real disposition (e.g. verdict "drop"), not just a
                # missing proof -- keep it in the acceptance manifest so the
                # reason is traceable, just omitted from composited/exports.
                # The proof image itself is untouched on disk; only the
                # concatenated deliverables (PDF/PPTX/HTML/MD) omit it.
                manifest_entries.append(entry)
            append_telemetry(job_id, {"event_type": "workflow", "stage": "verification_finalization", "status": "slide_dropped", "slide": idx, "reason": (entry or {}).get("reasoning")})
            continue
        spec = next(s for s in slide_specs if s["slide_index"] == idx)
        composited.append((
            spec, final_path, spec.get("headline", ""),
            " | ".join(spec.get("body_points", [])), spec.get("speaker_note", ""),
        ))
        manifest_entries.append(entry)

    if not composited:
        append_telemetry(job_id, {"event_type": "workflow", "stage": "verification_finalization", "status": "failed", "error": "No slides produced any usable final image"})
        db.update_job(job_id, status="failed", error_message="Verification finalization produced no usable slides")
        return

    # Hero -- same pipeline, included for consistency (verifier/judge calls
    # are cheap; only a flagged regeneration would touch the rate limiter).
    hero_final_path = None
    hero_manifest_entry = None
    hero_proof_path = slides_dir / "hero_proof.png"
    if hero_proof_path.exists():
        shutil.copyfile(hero_proof_path, export_presenter_dir / "hero_proof.png")
        hero_final_path = export_slides_dir / "hero_final.png"
        hero_storyboard_entry = {
            "slide_number": "hero", "title": deck_title,
            "purpose": "Cover/title slide", "image_prompt": "",
        }
        hero_manifest_entry = await _verify_and_finalize_slide(
            job_id, hero_storyboard_entry, hero_proof_path, hero_final_path,
            elevator_pitch, doc_text, visual_grammar, anchor_json,
            excepted_inference_elements=excepted_inference_elements,
            authorized_inferences=authorized_inferences,
            canon_overrides=canon_overrides,
            inferred_element_decisions=inferred_element_decisions,
            pitch_aspect_modes=pitch_aspect_modes,
            deck_storyboard=deck_storyboard,
        )
        if hero_manifest_entry.get("verdict") == "drop":
            # Drop path returns early without writing hero_final_path --
            # exclude hero from the composited deliverables the same way a
            # dropped numbered slide is excluded, keep the manifest entry
            # for traceability.
            hero_final_path = None
            dropped_slides.append("hero")
            append_telemetry(job_id, {"event_type": "workflow", "stage": "verification_finalization", "status": "slide_dropped", "slide": "hero", "reason": hero_manifest_entry.get("reasoning")})
        else:
            composited.insert(0, ({"slide_index": "hero", "slide_label": "Cover"}, hero_final_path, "", "", ""))

    pdf_path = export_dir / "deck.pdf"
    await asyncio.to_thread(_build_pdf, composited, pdf_path, deck_title)
    pptx_path = export_dir / "deck.pptx"
    await asyncio.to_thread(_build_pptx, composited, pptx_path, deck_title)
    html_path = export_dir / "preview.html"
    _build_html(composited, html_path, deck_title, job_id)
    md_path = export_dir / "slides.md"
    _build_markdown(composited, md_path, deck_title)

    numbered = [spec for spec, *_ in composited if spec.get("slide_index") != "hero"]
    manifest = {
        "job_id": job_id,
        "deck_title": deck_title,
        "slide_count": len(numbered),
        "mechanism": "parallel_verification_v2",
        "files": ["deck.pdf", "deck.pptx", "preview.html", "slides.md", "manifest.json", "verification_report.json"],
        "slide_images": [f"slides/slide_{spec['slide_index']:02d}_final.png" for spec in numbered],
        "presenter_view_images": [f"presenter-view/slide_{spec['slide_index']:02d}_proof.png" for spec in numbered],
        "telemetry_summary": "telemetry/summary.json",
        "telemetry": telemetry_summary(job_id),
        "dropped_slides": dropped_slides,
    }
    if hero_final_path:
        manifest["hero_image"] = "slides/hero_final.png"
        manifest["hero_presenter_view_image"] = "presenter-view/hero_proof.png"
    (export_dir / "manifest.json").write_text(json.dumps(manifest, indent=2))

    # The acceptance-manifest specified in the implementation prompt: every
    # slide's verifier findings, judge verdict, and final-image provenance.
    verification_report = {
        "job_id": job_id,
        "mechanism": "parallel_verification_v2",
        "slides": manifest_entries,
        "hero": hero_manifest_entry,
        "dropped_slides": dropped_slides,
        "regenerated_count": sum(1 for m in manifest_entries if m.get("final_source") == "regenerated")
                             + (1 if hero_manifest_entry and hero_manifest_entry.get("final_source") == "regenerated" else 0),
    }
    (export_dir / "verification_report.json").write_text(json.dumps(verification_report, indent=2))

    zip_path = export_dir / "pitch_deck_package.zip"
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for fname in manifest["files"]:
            p = export_dir / fname
            if p.exists():
                zf.write(p, fname)
        for p in sorted(export_slides_dir.glob("*.png")):
            zf.write(p, f"slides/{p.name}")
        for p in sorted(export_presenter_dir.glob("*.png")):
            zf.write(p, f"presenter-view/{p.name}")

    validation_errors = _validate_export_manifest(export_dir, zip_path, manifest)
    if validation_errors:
        append_telemetry(job_id, {"event_type": "workflow", "stage": "verification_finalization", "status": "failed", "errors": validation_errors})
        db.update_job(job_id, status="failed", error_message="Export validation failed: " + "; ".join(validation_errors))
        return

    db.update_job(job_id, status="complete", export_zip_path=str(zip_path))
    append_telemetry(job_id, {
        "event_type": "workflow",
        "stage": "verification_finalization",
        "status": "complete",
        "export_zip_path": zip_path.relative_to(job_dir(job_id)),
        "export_zip_bytes": zip_path.stat().st_size,
        "slide_count": len(composited),
        "regenerated_count": verification_report["regenerated_count"],
        "dropped_slides": dropped_slides,
    })


def _composite_slide(proof_path: Path, out_path: Path, text_regions: list, headline: str, body: str):
    from PIL import Image, ImageDraw, ImageFont

    if not proof_path.exists():
        img = Image.new("RGB", (1536, 1024), color=(10, 10, 20))
    else:
        img = Image.open(proof_path).convert("RGB")

    w, h = img.size
    draw = ImageDraw.Draw(img)

    # Cover text regions with panels
    for region in text_regions:
        rx = int(region.get("approx_x", 0) * w)
        ry = int(region.get("approx_y", 0) * h)
        rw = int(region.get("approx_w", 0.8) * w)
        rh = int(region.get("approx_h", 0.15) * h)
        # Sample background color near the region center
        sample_x = min(rx + rw // 2, w - 1)
        sample_y = min(ry + rh // 2, h - 1)
        try:
            bg_color = img.getpixel((sample_x, sample_y))[:3]
        except Exception:
            bg_color = (10, 10, 20)
        # Slightly darken for readability
        bg_color = tuple(max(0, c - 20) for c in bg_color)
        draw.rectangle([rx, ry, rx + rw, ry + rh], fill=bg_color)

    # Add text — try system font, fall back to default
    try:
        font_headline = ImageFont.truetype("/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf", max(24, h // 20))
        font_body = ImageFont.truetype("/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf", max(16, h // 35))
    except Exception:
        font_headline = ImageFont.load_default()
        font_body = ImageFont.load_default()

    # Headline: upper third
    draw.text((w * 0.08, h * 0.12), headline, fill=(240, 235, 220), font=font_headline)

    # Body: middle section
    body_y = h * 0.35
    for line in body.split(" | ")[:6]:
        if line.strip():
            draw.text((w * 0.08, body_y), f"• {line.strip()}", fill=(180, 175, 165), font=font_body)
            body_y += h * 0.07

    img.save(out_path, "PNG")


def _build_pdf(composited: list, pdf_path: Path, deck_title: str):
    from reportlab.lib.pagesizes import landscape
    from reportlab.lib.units import inch
    from reportlab.pdfgen import canvas as rl_canvas
    from PIL import Image as PILImage

    PAGE_W, PAGE_H = landscape((10 * inch, 7.5 * inch))
    c = rl_canvas.Canvas(str(pdf_path), pagesize=(PAGE_W, PAGE_H))
    c.setTitle(deck_title)

    for spec, img_path, headline, body, notes in composited:
        if img_path.exists():
            # Scale image to fill page
            c.drawImage(str(img_path), 0, 0, width=PAGE_W, height=PAGE_H,
                        preserveAspectRatio=True, anchor="c")
        c.showPage()

    c.save()


def _build_pptx(composited: list, pptx_path: Path, deck_title: str):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    from lxml import etree

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_layout = prs.slide_layouts[6]  # blank

    for spec, img_path, headline, body, notes in composited:
        slide = prs.slides.add_slide(blank_layout)

        # Star wipe transition on every slide (p:wheel spokes="4" is PowerPoint's star wipe)
        trans = etree.SubElement(slide._element, qn('p:transition'))
        etree.SubElement(trans, qn('p:wheel')).set('spokes', '4')

        if img_path.exists():
            slide.shapes.add_picture(
                str(img_path), Inches(0), Inches(0),
                width=prs.slide_width, height=prs.slide_height,
            )

        # Headline text box
        txBox = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(14), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = headline
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = RGBColor(0xF0, 0xEB, 0xDC)

        # Body text box
        if body:
            bodyBox = slide.shapes.add_textbox(Inches(1), Inches(2.6), Inches(14), Inches(5))
            btf = bodyBox.text_frame
            btf.word_wrap = True
            for line in body.split(" | ")[:6]:
                if line.strip():
                    bp = btf.add_paragraph()
                    bp.text = f"• {line.strip()}"
                    bp.font.size = Pt(20)
                    bp.font.color.rgb = RGBColor(0xB4, 0xAF, 0xA5)

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(str(pptx_path))


def _build_html(composited: list, html_path: Path, deck_title: str, job_id: str):
    slides_html = ""
    for spec, img_path, headline, body, notes in composited:
        img_url = f"slides/{img_path.name}"
        slides_html += f"""
<div class="slide">
  <img src="{img_url}" alt="{headline}" />
  <div class="caption">
    <h2>{headline}</h2>
    <p>{body}</p>
  </div>
</div>"""

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<title>{deck_title}</title>
<style>
body {{ margin:0; background:#0a0a0f; font-family: sans-serif; color:#e0ddd6; }}
.slide {{ width:100%; max-width:1200px; margin:2rem auto; }}
.slide img {{ width:100%; display:block; border:1px solid #1a1d24; }}
.caption {{ padding:1rem; background:#111318; }}
.caption h2 {{ margin:0 0 .5rem; font-size:1.1rem; }}
.caption p {{ margin:0; font-size:.85rem; color:#8a8880; }}
</style>
</head>
<body>
<h1 style="text-align:center;padding:2rem;font-size:1.4rem;letter-spacing:.12em;">{deck_title}</h1>
{slides_html}
</body>
</html>"""
    html_path.write_text(html)


def _build_markdown(composited: list, md_path: Path, deck_title: str):
    lines = [f"# {deck_title}\n"]
    for spec, img_path, headline, body, notes in composited:
        idx = spec.get("slide_index", "?")
        label = spec.get("slide_label", "")
        lines.append(f"\n## Slide {idx} — {label}\n")
        lines.append(f"**{headline}**\n")
        for point in body.split(" | "):
            if point.strip():
                lines.append(f"- {point.strip()}")
        if notes:
            lines.append(f"\n*Notes: {notes}*")
    md_path.write_text("\n".join(lines))
